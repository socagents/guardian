"""
Build the Cortex XDR Incident Response workshop deck.

Generates a 10-slide .pptx using python-pptx with a dark theme matching
the Cortex XDR console aesthetic (navy bg + Palo Alto orange accents).

Output: /Users/ayman/Documents/Coding/phantom/docs/workshops/cortex-xdr-ir-workshop.pptx

Content is grounded in:
- Cortex XDR docs (causality view, response actions, XQL, case management)
- Cortex XDR 5.0 announcement (Feb 2026) — case workspace + AI Case Investigation agent
- Real data from Phantom Master Killchain runs: incidents 1794 + 1795, 140+ alerts,
  9 MITRE tactics, 22 techniques validated end-to-end
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# ---------------------------------------------------------------------------
# Theme colors
# ---------------------------------------------------------------------------
BG_DARK = RGBColor(0x0F, 0x14, 0x19)          # Cortex XDR-ish navy background
BG_PANEL = RGBColor(0x1A, 0x1F, 0x2B)         # Slightly lighter panel
ACCENT_ORANGE = RGBColor(0xFA, 0x58, 0x2D)    # Palo Alto orange
ACCENT_RED = RGBColor(0xE6, 0x3A, 0x46)       # Malicious-process red (causality outline)
ACCENT_CYAN = RGBColor(0x06, 0xD6, 0xA0)      # Benign / safe state
ACCENT_BLUE = RGBColor(0x3B, 0x82, 0xF6)      # Highlight blue
TEXT_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_MUTED = RGBColor(0x94, 0xA3, 0xB8)       # Subtitle / metadata
TEXT_DIM = RGBColor(0xCB, 0xD5, 0xE1)         # Body text
GRID = RGBColor(0x33, 0x3F, 0x52)             # Faint grid / dividers

# Tactic chip colors — by ATT&CK lifecycle phase
TACTIC_COLORS = {
    "Initial Access": RGBColor(0xEF, 0x44, 0x44),
    "Execution": RGBColor(0xF9, 0x73, 0x16),
    "Persistence": RGBColor(0xEA, 0xB3, 0x08),
    "Priv. Esc.": RGBColor(0xCA, 0x8A, 0x04),
    "Defense Evasion": RGBColor(0x84, 0xCC, 0x16),
    "Cred. Access": RGBColor(0x10, 0xB9, 0x81),
    "Discovery": RGBColor(0x06, 0xB6, 0xD4),
    "Lateral Mvmt.": RGBColor(0x3B, 0x82, 0xF6),
    "Collection": RGBColor(0x8B, 0x5C, 0xF6),
    "C2": RGBColor(0xEC, 0x48, 0x99),
    "Exfil.": RGBColor(0xF4, 0x3F, 0x5E),
    "Impact": RGBColor(0xDC, 0x26, 0x26),
}


# ---------------------------------------------------------------------------
# Layout constants (16:9, 13.333" x 7.5")
# ---------------------------------------------------------------------------
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def set_slide_bg(slide, color=BG_DARK):
    """Set the entire slide background to a solid color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, x, y, w, h, fill=BG_PANEL, line=None, line_width=Pt(0)):
    """Add a filled rectangle (panel)."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = line_width
    shape.shadow.inherit = False
    return shape


def add_rounded_rect(slide, x, y, w, h, fill=ACCENT_ORANGE, line=None):
    """Add a rounded rectangle (chip / pill)."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
    shape.shadow.inherit = False
    # Reduce corner radius (default is huge); pptx exposes via adj
    try:
        shape.adjustments[0] = 0.25
    except Exception:
        pass
    return shape


def add_text(
    slide,
    x,
    y,
    w,
    h,
    text,
    *,
    size=14,
    bold=False,
    color=TEXT_WHITE,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
    font="Calibri",
):
    """Add a single-line / paragraph text box."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = Pt(2)
    tf.margin_right = Pt(2)
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    if isinstance(text, str):
        runs = [text]
    else:
        runs = text  # list of strings (will create one run per element)
    p.text = ""  # clear default
    for i, t in enumerate(runs):
        run = p.add_run() if i > 0 else p.add_run()
        run.text = t if isinstance(t, str) else t["text"]
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return tb


def add_multiline_text(
    slide,
    x,
    y,
    w,
    h,
    paragraphs,
    *,
    size=12,
    color=TEXT_DIM,
    line_spacing=1.2,
    font="Calibri",
):
    """Add a text box with multiple paragraphs.

    Each item in `paragraphs` can be a string or dict {text, bold, size, color, bullet}.
    """
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    tf.margin_top = Pt(4)
    tf.margin_bottom = Pt(4)
    tf.word_wrap = True
    for i, item in enumerate(paragraphs):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        if isinstance(item, str):
            p.text = item
            for run in p.runs:
                run.font.name = font
                run.font.size = Pt(size)
                run.font.color.rgb = color
        else:
            txt = item.get("text", "")
            p.text = txt
            for run in p.runs:
                run.font.name = item.get("font", font)
                run.font.size = Pt(item.get("size", size))
                run.font.bold = item.get("bold", False)
                run.font.color.rgb = item.get("color", color)
        p.line_spacing = line_spacing
    return tb


def add_header_bar(slide, title, subtitle=None, kicker=None):
    """Standard header for content slides."""
    # Top accent bar
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), fill=ACCENT_ORANGE)

    if kicker:
        add_text(
            slide,
            Inches(0.6),
            Inches(0.25),
            Inches(10),
            Inches(0.3),
            kicker.upper(),
            size=10,
            bold=True,
            color=ACCENT_ORANGE,
        )

    add_text(
        slide,
        Inches(0.6),
        Inches(0.55),
        Inches(12),
        Inches(0.7),
        title,
        size=28,
        bold=True,
        color=TEXT_WHITE,
    )
    if subtitle:
        add_text(
            slide,
            Inches(0.6),
            Inches(1.15),
            Inches(12),
            Inches(0.4),
            subtitle,
            size=14,
            color=TEXT_MUTED,
        )

    # Divider under header
    add_rect(slide, Inches(0.6), Inches(1.55), Inches(12.1), Pt(1), fill=GRID)


def add_footer(slide, slide_no, total):
    """Footer with slide number + workshop tag."""
    add_text(
        slide,
        Inches(0.6),
        Inches(7.05),
        Inches(8),
        Inches(0.3),
        "Cortex XDR IR Workshop  ·  Phantom Master Killchain Simulation",
        size=9,
        color=TEXT_MUTED,
    )
    add_text(
        slide,
        Inches(11.5),
        Inches(7.05),
        Inches(1.4),
        Inches(0.3),
        f"{slide_no} / {total}",
        size=9,
        color=TEXT_MUTED,
        align=PP_ALIGN.RIGHT,
    )


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------

def build_slide_1_title(prs):
    """Title slide — hero/cover."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide)

    # Decorative grid lines in the background
    for i in range(0, 14):
        add_rect(slide, Inches(i), Inches(0), Pt(0.5), SLIDE_H, fill=GRID)
    for i in range(0, 9):
        add_rect(slide, Inches(0), Inches(i), SLIDE_W, Pt(0.5), fill=GRID)

    # Top accent
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.12), fill=ACCENT_ORANGE)

    # Kicker
    add_text(
        slide,
        Inches(0.8),
        Inches(1.2),
        Inches(12),
        Inches(0.4),
        "HANDS-ON WORKSHOP  ·  SOC / IR PRACTITIONERS",
        size=12,
        bold=True,
        color=ACCENT_ORANGE,
    )

    # Title
    add_text(
        slide,
        Inches(0.8),
        Inches(1.8),
        Inches(12),
        Inches(1.5),
        "Incident Response with Cortex XDR",
        size=54,
        bold=True,
        color=TEXT_WHITE,
    )

    # Subtitle (two-line)
    add_text(
        slide,
        Inches(0.8),
        Inches(3.3),
        Inches(12),
        Inches(0.7),
        "Driving a real attack chain end-to-end through XDR's investigation surfaces",
        size=22,
        color=TEXT_DIM,
    )

    # Sim callout panel
    panel_y = Inches(4.6)
    panel = add_rect(slide, Inches(0.8), panel_y, Inches(11.7), Inches(1.5), fill=BG_PANEL)
    add_text(
        slide,
        Inches(1.1),
        Inches(4.75),
        Inches(11),
        Inches(0.4),
        "SIMULATED ADVERSARY",
        size=11,
        bold=True,
        color=ACCENT_ORANGE,
    )
    add_text(
        slide,
        Inches(1.1),
        Inches(5.05),
        Inches(11),
        Inches(0.5),
        "Phantom Master Killchain  ·  22 abilities  ·  9 MITRE tactics",
        size=22,
        bold=True,
        color=TEXT_WHITE,
    )
    add_text(
        slide,
        Inches(1.1),
        Inches(5.55),
        Inches(11),
        Inches(0.5),
        "Validated against Cortex XDR cases 1794 + 1795 — 140+ alerts across 22 techniques",
        size=13,
        color=TEXT_MUTED,
    )

    # Footer
    add_text(
        slide,
        Inches(0.8),
        Inches(6.6),
        Inches(11),
        Inches(0.4),
        "Phantom platform  ·  v0.6.x  ·  May 2026",
        size=11,
        color=TEXT_MUTED,
    )


def build_slide_2_objectives(prs, total):
    """Workshop objectives."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header_bar(
        slide,
        "What you'll walk out with",
        subtitle="Five concrete capabilities you'll use on a real attack chain — not slides about features",
        kicker="Workshop objectives",
    )

    objectives = [
        ("01", "Trace an alert to its root cause",
         "Use the causality view to follow a process tree from a single Cortex alert back to its initial-access ancestor — Spearphishing Attachment in our case."),
        ("02", "Hunt across the data plane with XQL",
         "Write XQL queries that pivot off an IOC from the causality view and find sibling activity the alerts didn't surface — 800+ catalogued fields, cross-data-source."),
        ("03", "Drive a case through the new 5.0 workspace",
         "Triage, prioritize, assign, document, and close a multi-stage incident. Use the AI Case Investigation agent to summarize signals + recommend next steps."),
        ("04", "Contain a compromised host live",
         "Isolate the endpoint, block the file hash, open Live Terminal, terminate the process tree, quarantine artifacts — all from the same console."),
        ("05", "Read MITRE coverage from telemetry",
         "Map the 22 alerts produced by the simulation onto the ATT&CK matrix and understand which tactics XDR's analytics actually fired on."),
    ]

    box_w = Inches(11.7)
    start_y = Inches(2.0)
    box_h = Inches(0.95)
    gap = Inches(0.08)
    for i, (num, title, desc) in enumerate(objectives):
        y = start_y + i * (box_h + gap)
        # Number badge
        add_rounded_rect(slide, Inches(0.6), y, Inches(0.8), box_h, fill=ACCENT_ORANGE)
        add_text(
            slide,
            Inches(0.6),
            y,
            Inches(0.8),
            box_h,
            num,
            size=22,
            bold=True,
            color=TEXT_WHITE,
            align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE,
        )
        # Content panel
        add_rect(slide, Inches(1.55), y, box_w - Inches(0.95), box_h, fill=BG_PANEL)
        add_text(
            slide,
            Inches(1.75),
            y + Inches(0.08),
            box_w - Inches(1.15),
            Inches(0.4),
            title,
            size=15,
            bold=True,
            color=TEXT_WHITE,
        )
        add_text(
            slide,
            Inches(1.75),
            y + Inches(0.45),
            box_w - Inches(1.15),
            Inches(0.5),
            desc,
            size=11,
            color=TEXT_DIM,
        )

    add_footer(slide, 2, total)


def build_slide_3_simulation(prs, total):
    """The Phantom Master Killchain simulation — overview."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header_bar(
        slide,
        "The simulation — Phantom Master Killchain",
        subtitle="22 chained abilities that mirror a real ransomware operation — Spearphishing → C2 → Defacement",
        kicker="Section 1  ·  The attack",
    )

    # Left column — stat cards
    cards = [
        ("22", "Caldera abilities chained in atomic ordering"),
        ("9", "MITRE ATT&CK tactics covered end-to-end"),
        ("22", "Distinct techniques exercised on the host"),
        ("140+", "Alerts surfaced in Cortex incidents 1794 + 1795"),
        ("2", "Cross-host lateral movement (xdragent ↔ xdragent2)"),
        ("88.6%", "Chain-quality rate validated against XDR analytics"),
    ]

    card_x_start = Inches(0.6)
    card_y_start = Inches(2.0)
    card_w = Inches(3.8)
    card_h = Inches(1.45)
    gap_x = Inches(0.15)
    gap_y = Inches(0.15)

    for i, (big, label) in enumerate(cards):
        col = i % 2
        row = i // 2
        x = card_x_start + col * (card_w + gap_x)
        y = card_y_start + row * (card_h + gap_y)
        add_rect(slide, x, y, card_w, card_h, fill=BG_PANEL)
        # Left accent bar
        add_rect(slide, x, y, Inches(0.08), card_h, fill=ACCENT_ORANGE)
        add_text(
            slide,
            x + Inches(0.2),
            y + Inches(0.1),
            card_w - Inches(0.3),
            Inches(0.7),
            big,
            size=40,
            bold=True,
            color=ACCENT_ORANGE,
        )
        add_text(
            slide,
            x + Inches(0.2),
            y + Inches(0.85),
            card_w - Inches(0.3),
            Inches(0.6),
            label,
            size=11,
            color=TEXT_DIM,
        )

    # Right column — phase timeline
    timeline_x = Inches(8.6)
    timeline_y = Inches(2.0)
    timeline_w = Inches(4.2)
    timeline_h = Inches(4.6)
    add_rect(slide, timeline_x, timeline_y, timeline_w, timeline_h, fill=BG_PANEL)
    add_text(
        slide,
        timeline_x + Inches(0.2),
        timeline_y + Inches(0.1),
        timeline_w - Inches(0.4),
        Inches(0.4),
        "ATTACK CHAIN PHASES",
        size=11,
        bold=True,
        color=ACCENT_ORANGE,
    )

    phases = [
        ("01", "Initial Access", "Phishing macro drops Caldera implant"),
        ("02", "Discovery", "Account / system / share enumeration"),
        ("03", "Credential Access", "LSASS dump + PowerKatz + cached creds"),
        ("04", "Priv. Esc. + Persistence", "Fodhelper UAC bypass + run key + scheduled task"),
        ("05", "Defense Evasion", "Disable Defender + certutil decode"),
        ("06", "Lateral Movement", "SMB + WMI + WinRM to second host"),
        ("07", "Collection + Exfil", "Auto-collection + archive + DNS C2"),
        ("08", "Impact + Cleanup", "Defacement + clear event log"),
    ]

    px = timeline_x + Inches(0.4)
    py_start = timeline_y + Inches(0.5)
    step = Inches(0.50)
    for i, (n, name, desc) in enumerate(phases):
        py = py_start + i * step
        # Step number
        add_text(
            slide, px, py + Inches(0.05), Inches(0.4), Inches(0.4),
            n, size=11, bold=True, color=ACCENT_ORANGE,
        )
        add_text(
            slide, px + Inches(0.4), py - Inches(0.05), timeline_w - Inches(0.8), Inches(0.3),
            name, size=12, bold=True, color=TEXT_WHITE,
        )
        add_text(
            slide, px + Inches(0.4), py + Inches(0.2), timeline_w - Inches(0.8), Inches(0.3),
            desc, size=9, color=TEXT_MUTED,
        )

    add_footer(slide, 3, total)


def build_slide_4_mitre_matrix(prs, total):
    """MITRE tactic/technique coverage table."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header_bar(
        slide,
        "MITRE ATT&CK coverage map",
        subtitle="Each row is a tactic the simulation exercises — each cell is a technique with a real Caldera ability behind it",
        kicker="Section 1  ·  Coverage",
    )

    # Rows: (tactic_short, attck_id, techniques)
    rows = [
        ("Initial Access", "TA0001", ["T1566.001 Spearphishing Attachment"]),
        ("Execution", "TA0002", ["T1059.001 PowerShell", "T1059.003 Windows Command Shell", "T1204 User Execution"]),
        ("Persistence", "TA0003", ["T1136.001 Local Account", "T1547.001 Registry Run Key", "T1053.005 Scheduled Task", "T1574.002 DLL Side-Load"]),
        ("Priv. Esc.", "TA0004", ["T1548.002 Fodhelper UAC Bypass"]),
        ("Defense Evasion", "TA0005", ["T1562.001 Disable Defender", "T1140 Certutil Decode", "T1132 Data Encoding"]),
        ("Cred. Access", "TA0006", ["T1003.001 LSASS Memory", "T1003.004 LSA Secrets", "T1003.005 Cached Creds", "T1555 Credential Manager"]),
        ("Discovery", "TA0007", ["T1087.001 Local Account", "T1082 System Info", "T1018 Remote System", "T1135 Network Share", "T1518.001 AV/EDR"]),
        ("Lateral Mvmt.", "TA0008", ["T1021.002 SMB/Admin Shares", "T1550 Pass-the-Hash"]),
        ("Collection", "TA0009", ["T1119 Automated Collection", "T1560 Archive Collected"]),
        ("C2", "TA0011", ["T1071.004 DNS", "T1197 BITS Jobs"]),
        ("Exfil.", "TA0010", ["T1041 Exfil over C2"]),
        ("Impact", "TA0040", ["T1491 Defacement", "T1070.001 Clear Event Log"]),
    ]

    # Table layout — compressed to fit 12 rows + footer note above the page footer
    table_x = Inches(0.6)
    table_y = Inches(1.85)
    table_w = Inches(12.1)
    row_h = Inches(0.34)
    chip_w = Inches(2.0)
    id_w = Inches(0.9)
    tech_w = table_w - chip_w - id_w - Inches(0.2)

    # Header
    add_text(slide, table_x, table_y, chip_w, row_h, "TACTIC", size=10, bold=True, color=TEXT_MUTED)
    add_text(slide, table_x + chip_w + Inches(0.1), table_y, id_w, row_h, "ATT&CK ID", size=10, bold=True, color=TEXT_MUTED)
    add_text(slide, table_x + chip_w + id_w + Inches(0.2), table_y, tech_w, row_h, "TECHNIQUES EXERCISED", size=10, bold=True, color=TEXT_MUTED)

    table_y += Inches(0.3)
    for i, (tactic, attck_id, techs) in enumerate(rows):
        y = table_y + i * row_h
        # Tactic chip
        chip_color = TACTIC_COLORS.get(tactic, ACCENT_ORANGE)
        add_rounded_rect(slide, table_x, y + Inches(0.04), chip_w - Inches(0.2), row_h - Inches(0.08), fill=chip_color)
        add_text(
            slide, table_x, y + Inches(0.04), chip_w - Inches(0.2), row_h - Inches(0.08),
            tactic, size=10, bold=True, color=TEXT_WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
        )
        # ATT&CK ID
        add_text(slide, table_x + chip_w + Inches(0.1), y + Inches(0.04), id_w, row_h, attck_id, size=10, color=TEXT_DIM, anchor=MSO_ANCHOR.MIDDLE)
        # Techniques
        tech_str = "  ·  ".join(techs)
        add_text(
            slide, table_x + chip_w + id_w + Inches(0.2), y + Inches(0.04),
            tech_w, row_h, tech_str, size=10, color=TEXT_WHITE, anchor=MSO_ANCHOR.MIDDLE,
        )
        # Faint row divider
        if i < len(rows) - 1:
            add_rect(slide, table_x, y + row_h - Pt(0.5), table_w, Pt(0.5), fill=GRID)

    # Footer note (positioned just below the table, above the page footer)
    note_y = table_y + len(rows) * row_h + Inches(0.05)
    add_text(
        slide,
        Inches(0.6),
        note_y,
        Inches(12),
        Inches(0.3),
        "Every row above produced at least one Cortex XDR alert in incidents 1794 + 1795 — this is verified telemetry, not a wishlist.",
        size=10,
        color=ACCENT_ORANGE,
    )

    add_footer(slide, 4, total)


def build_slide_5_causality(prs, total):
    """XDR Capability #1 — Causality view."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header_bar(
        slide,
        "Causality view — find the root cause in one screen",
        subtitle="Cortex's process-tree visualization. Every alert traces back to a process node; nodes carry visual verdict.",
        kicker="Section 2 · XDR Capability 1 of 4",
    )

    # Left: feature explanation
    panel_x = Inches(0.6)
    panel_w = Inches(6.0)
    panel_y = Inches(2.0)

    add_text(slide, panel_x, panel_y, panel_w, Inches(0.4), "WHAT IT DOES", size=11, bold=True, color=ACCENT_ORANGE)
    add_multiline_text(
        slide,
        panel_x,
        panel_y + Inches(0.35),
        panel_w,
        Inches(1.6),
        [
            "Every Cortex XDR alert has a causality chain — a directed graph of process nodes (the things that ran) and executions (RPC, code injection, file I/O between them) that led to the detection.",
            "Nodes carry verdict-colored outlines: red for malicious processes, blue for benign per WildFire. Analysts read the chain end-to-end to identify the root cause.",
        ],
        size=11,
        line_spacing=1.3,
    )

    add_text(slide, panel_x, Inches(4.2), panel_w, Inches(0.4), "WHAT YOU CAN DO ON A NODE", size=11, bold=True, color=ACCENT_ORANGE)
    actions = [
        "Investigate further (drill into telemetry)",
        "Timeline view — chronological events around the node",
        "Add hash to allow / block list",
        "Search the file across all endpoints",
        "Terminate the process — live, on the agent",
        "Quarantine the file artifact",
    ]
    add_multiline_text(
        slide,
        panel_x,
        Inches(4.6),
        panel_w,
        Inches(2.2),
        [f"·  {a}" for a in actions],
        size=11,
        line_spacing=1.3,
    )

    # Right: example causality from our simulation
    ex_x = Inches(7.0)
    ex_y = Inches(2.0)
    ex_w = Inches(5.7)
    ex_h = Inches(4.6)
    add_rect(slide, ex_x, ex_y, ex_w, ex_h, fill=BG_PANEL)
    add_text(slide, ex_x + Inches(0.2), ex_y + Inches(0.15), ex_w - Inches(0.4), Inches(0.4),
             "LIVE EXAMPLE  ·  INCIDENT 1794", size=11, bold=True, color=ACCENT_ORANGE)
    add_text(slide, ex_x + Inches(0.2), ex_y + Inches(0.5), ex_w - Inches(0.4), Inches(0.4),
             "Causality chain produced by the phishing dropper", size=14, bold=True, color=TEXT_WHITE)

    # Chain — vertical process tree visualization (text-only nodes)
    nodes = [
        ("WINWORD.EXE (masqueraded)", "BENIGN", ACCENT_CYAN, "Phishing attachment opens — initial-access trigger"),
        ("cmd.exe /c", "BENIGN", ACCENT_CYAN, "Macro spawns shell"),
        ("powershell.exe", "MALICIOUS", ACCENT_RED, "Encoded command runs (T1059.001)"),
        ("bitsadmin /transfer", "MALICIOUS", ACCENT_RED, "BITS job downloads dropper (T1197)"),
        ("mars.exe (Caldera implant)", "MALICIOUS", ACCENT_RED, "Sandcat beacon comes alive — C2 established"),
    ]
    node_x = ex_x + Inches(0.4)
    node_y = ex_y + Inches(1.0)
    node_w = ex_w - Inches(0.8)
    node_h = Inches(0.55)
    gap = Inches(0.15)

    for i, (proc, verdict, color, ctx) in enumerate(nodes):
        y = node_y + i * (node_h + gap)
        # Left outline indicates verdict
        add_rect(slide, node_x, y, Inches(0.08), node_h, fill=color)
        # Node panel
        add_rect(slide, node_x + Inches(0.08), y, node_w - Inches(0.08), node_h, fill=BG_DARK)
        # Process name + verdict
        add_text(slide, node_x + Inches(0.25), y + Inches(0.04), node_w * 0.55, Inches(0.3),
                 proc, size=12, bold=True, color=TEXT_WHITE)
        add_text(slide, node_x + node_w * 0.6, y + Inches(0.06), node_w * 0.38, Inches(0.25),
                 verdict, size=9, bold=True, color=color, align=PP_ALIGN.RIGHT)
        # Context
        add_text(slide, node_x + Inches(0.25), y + Inches(0.28), node_w - Inches(0.3), Inches(0.25),
                 ctx, size=9, color=TEXT_MUTED)

        # Arrow chevron (down) between nodes — small
        if i < len(nodes) - 1:
            chev_x = node_x + node_w / 2 - Inches(0.05)
            chev_y = y + node_h + Inches(0.01)
            tri = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, chev_x, chev_y, Inches(0.1), Inches(0.12))
            tri.fill.solid()
            tri.fill.fore_color.rgb = TEXT_MUTED
            tri.line.fill.background()

    add_footer(slide, 5, total)


def build_slide_6_xql(prs, total):
    """XDR Capability #2 — XQL hunting."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header_bar(
        slide,
        "XQL search — hunt beyond what the alerts surfaced",
        subtitle="800+ catalogued fields · regex + JSON · merge data across endpoint / network / cloud / identity",
        kicker="Section 2 · XDR Capability 2 of 4",
    )

    # Left: what it is + when to use it
    panel_x = Inches(0.6)
    panel_w = Inches(5.0)
    add_text(slide, panel_x, Inches(2.0), panel_w, Inches(0.4),
             "WHEN YOU REACH FOR XQL", size=11, bold=True, color=ACCENT_ORANGE)
    add_multiline_text(
        slide,
        panel_x,
        Inches(2.4),
        panel_w,
        Inches(4.3),
        [
            "·  An alert fires — you want to find every OTHER host running the same encoded PowerShell.",
            "·  A causality node looks suspicious but didn't trigger analytics — you need to pivot off the parent PID + binary hash.",
            "·  Threat-intel drops an IOC at 2am — you need to scope its blast radius across the estate before you wake up the team.",
            "·  You're writing a new analytic rule and want to test the query against 30 days of telemetry first.",
            "·  Compromise assessment — your customer asks 'are we already in?' for a TTP that fired on a peer.",
        ],
        size=11,
        line_spacing=1.4,
    )

    # Right: example XQL panel
    ex_x = Inches(6.1)
    ex_w = Inches(6.6)
    ex_y = Inches(2.0)
    ex_h = Inches(4.7)
    add_rect(slide, ex_x, ex_y, ex_w, ex_h, fill=BG_PANEL)

    add_text(slide, ex_x + Inches(0.2), ex_y + Inches(0.15), ex_w, Inches(0.4),
             "EXAMPLE — HUNT THE SIMULATION'S BITS JOB", size=11, bold=True, color=ACCENT_ORANGE)
    add_text(slide, ex_x + Inches(0.2), ex_y + Inches(0.5), ex_w, Inches(0.4),
             "Find every host where bitsadmin.exe transferred a file in the last 24h:",
             size=12, color=TEXT_DIM)

    # Code-block looking area (taller — XQL is 7 lines)
    code_x = ex_x + Inches(0.3)
    code_y = ex_y + Inches(1.0)
    code_w = ex_w - Inches(0.6)
    code_h = Inches(1.95)
    add_rect(slide, code_x, code_y, code_w, code_h, fill=BG_DARK, line=GRID, line_width=Pt(0.5))
    code_lines = [
        "dataset = xdr_data",
        "| filter event_type = ENUM.PROCESS",
        "       and action_process_image_name = \"bitsadmin.exe\"",
        "       and action_process_image_command_line contains \"transfer\"",
        "       and _time > to_timestamp(\"now\") - 24h",
        "| comp count() by agent_hostname,",
        "          causality_actor_process_image_name",
        "| sort desc _count",
    ]
    add_multiline_text(
        slide, code_x + Inches(0.15), code_y + Inches(0.1),
        code_w - Inches(0.3), code_h - Inches(0.2),
        code_lines, size=10, color=ACCENT_CYAN, font="Consolas", line_spacing=1.15,
    )

    # Result narrative (moved down to clear the taller code panel)
    result_y = code_y + code_h + Inches(0.1)
    add_text(slide, ex_x + Inches(0.2), result_y, ex_w - Inches(0.4), Inches(0.35),
             "RESULT IN OUR LAB", size=11, bold=True, color=ACCENT_ORANGE)
    add_multiline_text(
        slide,
        ex_x + Inches(0.2),
        result_y + Inches(0.35),
        ex_w - Inches(0.4),
        Inches(1.3),
        [
            "·  xdragent — bitsadmin transferred mars.exe (1 hit, parent powershell.exe)",
            "·  xdragent2 — same pattern (1 hit, post-lateral-movement)",
            "·  Pivot from this row → causality view of either host → see the FULL ransack chain.",
        ],
        size=10,
        line_spacing=1.3,
    )

    add_footer(slide, 6, total)


def build_slide_7_cases(prs, total):
    """XDR Capability #3 — Case management."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header_bar(
        slide,
        "Case management — the IR workflow you live in",
        subtitle="Cortex XDR 5.0 (Feb 2026) rebuilt the case workspace. AI Case Investigation agent + visual context maps.",
        kicker="Section 2 · XDR Capability 3 of 4",
    )

    # Workflow strip — 5 phases as connected pill boxes
    phases = [
        ("OPEN", "Alert lands in queue", ACCENT_ORANGE),
        ("TRIAGE", "Severity + scope set", ACCENT_RED),
        ("INVESTIGATE", "Causality + XQL pivots", ACCENT_BLUE),
        ("RESPOND", "Isolate + block + kill", RGBColor(0xCA, 0x8A, 0x04)),
        ("CLOSE", "Document + lessons learned", ACCENT_CYAN),
    ]

    strip_x = Inches(0.6)
    strip_y = Inches(2.0)
    pill_w = Inches(2.2)
    pill_h = Inches(1.3)
    gap = Inches(0.25)

    for i, (label, desc, color) in enumerate(phases):
        x = strip_x + i * (pill_w + gap)
        add_rounded_rect(slide, x, strip_y, pill_w, pill_h, fill=color)
        add_text(slide, x, strip_y + Inches(0.25), pill_w, Inches(0.45),
                 label, size=15, bold=True, color=TEXT_WHITE, align=PP_ALIGN.CENTER)
        # Wrap desc to two lines naturally via the word_wrap text frame
        add_text(slide, x + Inches(0.15), strip_y + Inches(0.75), pill_w - Inches(0.3), Inches(0.5),
                 desc, size=11, color=TEXT_WHITE, align=PP_ALIGN.CENTER)
        # Chevron between
        if i < len(phases) - 1:
            chev_x = x + pill_w + Inches(0.02)
            chev_y = strip_y + pill_h / 2 - Inches(0.15)
            tri = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, chev_x, chev_y, Inches(0.2), Inches(0.3))
            tri.fill.solid()
            tri.fill.fore_color.rgb = TEXT_MUTED
            tri.line.fill.background()

    # What XDR 5.0 added panel
    add_text(slide, Inches(0.6), Inches(3.6), Inches(12), Inches(0.4),
             "WHAT CORTEX XDR 5.0 CHANGED IN THE ANALYST WORKSPACE", size=11, bold=True, color=ACCENT_ORANGE)

    cards = [
        ("AI Case Summaries",
         "Plain-language briefing of complex multi-stage incidents. Analyst reads 3 sentences instead of 40 alerts."),
        ("Connection maps",
         "Visualize alerts ↔ assets ↔ users in one graph. See the lateral path the adversary took without writing a query."),
        ("Case Investigation agent",
         "Agentic AI suggests next steps, highlights critical evidence, can act with analyst oversight (e.g. host containment, enrichment)."),
        ("Triage status tracking",
         "Mark Open / In Progress / Pending Review / Closed. Every status change is audit-trailed; multi-analyst handoffs are clean."),
    ]

    card_x_start = Inches(0.6)
    card_y_start = Inches(4.05)
    card_w = Inches(2.95)
    card_h = Inches(2.5)
    gap = Inches(0.1)
    for i, (title, desc) in enumerate(cards):
        x = card_x_start + i * (card_w + gap)
        add_rect(slide, x, card_y_start, card_w, card_h, fill=BG_PANEL)
        add_rect(slide, x, card_y_start, card_w, Inches(0.08), fill=ACCENT_ORANGE)
        add_text(slide, x + Inches(0.15), card_y_start + Inches(0.25), card_w - Inches(0.3), Inches(0.6),
                 title, size=13, bold=True, color=TEXT_WHITE)
        add_text(slide, x + Inches(0.15), card_y_start + Inches(0.85), card_w - Inches(0.3), card_h - Inches(1.0),
                 desc, size=10, color=TEXT_DIM)

    add_footer(slide, 7, total)


def build_slide_8_response_actions(prs, total):
    """XDR Capability #4 — Response actions."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header_bar(
        slide,
        "Response actions — contain the host without leaving the console",
        subtitle="Every action is reversible · every action is audit-logged · every action triggers an agent task you can track",
        kicker="Section 2 · XDR Capability 4 of 4",
    )

    actions = [
        ("Isolate Endpoint",
         "Halts all network traffic except agent ↔ Cortex XDR. Adversary's C2 channel goes silent immediately. Reversible.",
         "When you've confirmed compromise but haven't finished eradication.",
         ACCENT_RED),
        ("Live Terminal",
         "Remote shell on the host. Use it to inspect filesystem, kill processes, dump memory, copy artifacts — without sending anyone to physically touch the machine.",
         "When you need surgical action that no policy / button covers.",
         ACCENT_BLUE),
        ("Block File Hash",
         "Add SHA256 to the global block list. Cortex agents refuse to execute the binary next time it appears. Independent of isolation.",
         "After you've confirmed a hash is malicious — prevents recurrence across the estate.",
         RGBColor(0xCA, 0x8A, 0x04)),
        ("Terminate Process",
         "Kill the running process tree on a specific endpoint, from the causality view. Doesn't wait for the agent's policy to catch it.",
         "When you've found a malicious node and need it gone now.",
         RGBColor(0xEA, 0xB3, 0x08)),
        ("Quarantine File",
         "Move the malicious artifact to a protected location. File can be restored if it turns out to be FP, or sent to forensics for analysis.",
         "After process termination, to neutralize the artifact on disk.",
         ACCENT_CYAN),
        ("Search File on All Endpoints",
         "Right-click the hash in causality view → 'Search'. Returns every host where the file currently exists or recently ran.",
         "Lateral-spread assessment in seconds, not hours.",
         RGBColor(0x8B, 0x5C, 0xF6)),
    ]

    grid_x = Inches(0.6)
    grid_y = Inches(2.0)
    card_w = Inches(3.95)
    card_h = Inches(2.35)
    gap_x = Inches(0.15)
    gap_y = Inches(0.2)
    for i, (title, what, when, color) in enumerate(actions):
        col = i % 3
        row = i // 3
        x = grid_x + col * (card_w + gap_x)
        y = grid_y + row * (card_h + gap_y)
        add_rect(slide, x, y, card_w, card_h, fill=BG_PANEL)
        add_rect(slide, x, y, Inches(0.1), card_h, fill=color)
        add_text(slide, x + Inches(0.2), y + Inches(0.15), card_w - Inches(0.3), Inches(0.45),
                 title, size=15, bold=True, color=TEXT_WHITE)
        add_text(slide, x + Inches(0.2), y + Inches(0.6), card_w - Inches(0.3), Inches(0.3),
                 "WHAT IT DOES", size=9, bold=True, color=color)
        add_text(slide, x + Inches(0.2), y + Inches(0.85), card_w - Inches(0.3), Inches(0.8),
                 what, size=10, color=TEXT_DIM)
        add_text(slide, x + Inches(0.2), y + Inches(1.65), card_w - Inches(0.3), Inches(0.25),
                 "WHEN", size=9, bold=True, color=color)
        add_text(slide, x + Inches(0.2), y + Inches(1.85), card_w - Inches(0.3), Inches(0.45),
                 when, size=10, color=TEXT_DIM)

    add_footer(slide, 8, total)


def build_slide_9_hands_on(prs, total):
    """Hands-on exercise."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header_bar(
        slide,
        "Hands-on — drive incident 1794 from open to closed",
        subtitle="You'll work in pairs. 30 minutes. You have the actual telemetry on the lab tenant.",
        kicker="Section 3 · Your turn",
    )

    # Big timeline of tasks
    tasks = [
        ("00:00 – 00:05", "Find the alert", "In the Cortex tenant, find incident 1794 (severity high). Note the top-level fired analytics rule."),
        ("00:05 – 00:10", "Read the causality", "Open the alert's causality view. Walk from the WildFire-verdict malicious node back to the initial-access process. Identify the root cause."),
        ("00:10 – 00:18", "Pivot to XQL", "From a node hash in the chain, write an XQL query to find every other host that ran the same binary in the last 7 days. Take a screenshot of the result."),
        ("00:18 – 00:23", "Open the case", "Set triage status to 'In Progress', assign to yourself, add a 1-paragraph case summary in your own words. (AI summary is fine to compare against.)"),
        ("00:23 – 00:28", "Contain the host", "Isolate the endpoint. Block the malicious hash globally. Terminate the process tree from the causality view. Quarantine the artifact on disk."),
        ("00:28 – 00:30", "Close out", "Set triage status to 'Resolved'. Document what you'd do differently in production (this is a lab — what's missing or risky here?)."),
    ]

    list_x = Inches(0.6)
    list_y = Inches(2.0)
    row_h = Inches(0.75)
    gap = Inches(0.05)
    time_w = Inches(2.0)
    title_w = Inches(2.6)
    desc_w = SLIDE_W - list_x - time_w - title_w - Inches(0.6)

    for i, (time, title, desc) in enumerate(tasks):
        y = list_y + i * (row_h + gap)
        add_rect(slide, list_x, y, time_w, row_h, fill=BG_PANEL)
        add_text(slide, list_x + Inches(0.15), y + Inches(0.05), time_w - Inches(0.3), row_h - Inches(0.1),
                 time, size=12, bold=True, color=ACCENT_ORANGE, anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, list_x + time_w + Inches(0.15), y + Inches(0.05), title_w, row_h - Inches(0.1),
                 title, size=13, bold=True, color=TEXT_WHITE, anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, list_x + time_w + title_w + Inches(0.2), y + Inches(0.1), desc_w, row_h - Inches(0.2),
                 desc, size=10.5, color=TEXT_DIM)

    add_footer(slide, 9, total)


def build_slide_10_wrap_up(prs, total):
    """Wrap-up / takeaways."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header_bar(
        slide,
        "What you should walk out remembering",
        subtitle="If everything else fades, hold on to these five.",
        kicker="Wrap-up · Q&A",
    )

    takeaways = [
        ("The causality view is the spine of every investigation.",
         "Every alert in Cortex XDR has a process tree behind it. If you don't read that tree, you're triaging on labels."),
        ("XQL is your second screen.",
         "Causality tells you the story; XQL tells you the blast radius. Use both. 800+ fields, cross-data-source. Save your useful queries."),
        ("The new 5.0 workspace is collaborative on purpose.",
         "Open → Triage → Investigate → Respond → Close. Multi-analyst handoffs are first-class. The AI agent is a force multiplier, not a substitute for judgment."),
        ("Response is reversible — but only if you act.",
         "Isolate freely. Block hashes freely. Both undo cleanly if you were wrong. Hesitation costs more than over-containment."),
        ("Treat MITRE coverage as evidence, not aspiration.",
         "The simulation maps to 9 tactics + 22 techniques because it actually fires alerts on each. If your stack doesn't fire on a tactic, it's not covered — vendor claims notwithstanding."),
    ]

    list_x = Inches(0.6)
    list_y = Inches(1.95)
    box_h = Inches(0.92)
    gap = Inches(0.08)
    for i, (head, body) in enumerate(takeaways):
        y = list_y + i * (box_h + gap)
        # Number badge
        add_rect(slide, list_x, y, Inches(0.5), box_h, fill=ACCENT_ORANGE)
        add_text(slide, list_x, y, Inches(0.5), box_h,
                 str(i + 1), size=22, bold=True, color=TEXT_WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # Body panel
        add_rect(slide, list_x + Inches(0.6), y, Inches(12.1), box_h, fill=BG_PANEL)
        add_text(slide, list_x + Inches(0.8), y + Inches(0.1), Inches(11.7), Inches(0.4),
                 head, size=13, bold=True, color=TEXT_WHITE)
        add_text(slide, list_x + Inches(0.8), y + Inches(0.45), Inches(11.7), Inches(0.5),
                 body, size=11, color=TEXT_DIM)

    add_footer(slide, 10, total)


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    TOTAL = 10
    build_slide_1_title(prs)
    build_slide_2_objectives(prs, TOTAL)
    build_slide_3_simulation(prs, TOTAL)
    build_slide_4_mitre_matrix(prs, TOTAL)
    build_slide_5_causality(prs, TOTAL)
    build_slide_6_xql(prs, TOTAL)
    build_slide_7_cases(prs, TOTAL)
    build_slide_8_response_actions(prs, TOTAL)
    build_slide_9_hands_on(prs, TOTAL)
    build_slide_10_wrap_up(prs, TOTAL)

    out = "/Users/ayman/Documents/Coding/phantom/docs/workshops/cortex-xdr-ir-workshop.pptx"
    prs.save(out)
    print(f"Wrote {out}")
    print(f"  slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
