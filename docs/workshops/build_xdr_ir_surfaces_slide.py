"""
Single-slide build matching the Cortex XDR marketing-slide style
(black bg, Palo Alto green accents, central white badge, 5 capability
cards in a row, bottom 'data sources' strip with green up-arrows,
right sidebar with three value-prop blocks separated by horizontal lines).

Content is the IR-capabilities map we sketched for the SOC analyst
workshop — surfaces inside Cortex XDR for working an incident case.

Output: /Users/ayman/Documents/Coding/phantom/docs/workshops/cortex-xdr-ir-surfaces.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# --- Style palette (matched against the reference screenshot) ---------------
BG_BLACK = RGBColor(0x00, 0x00, 0x00)
PALO_GREEN = RGBColor(0x00, 0xC1, 0x6B)
GREEN_DIM = RGBColor(0x00, 0x80, 0x4A)
GREEN_GLOW = RGBColor(0x02, 0x8A, 0x4D)
CARD_DARK = RGBColor(0x12, 0x14, 0x18)
CARD_OUTLINE = RGBColor(0x33, 0x39, 0x40)
CLOUD_GRAY = RGBColor(0x1F, 0x25, 0x2E)
CLOUD_DARK = RGBColor(0x12, 0x17, 0x1E)
TEXT_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_DIM = RGBColor(0xCB, 0xD5, 0xE1)
TEXT_MUTED = RGBColor(0x94, 0xA3, 0xB8)
DIVIDER = RGBColor(0x00, 0x80, 0x4A)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# --- Helpers ----------------------------------------------------------------
def set_bg(slide, color=BG_BLACK):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, x, y, w, h, fill=CARD_DARK, line=None, line_width=Pt(0), rounded=False, radius=0.08):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE, x, y, w, h
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = line_width
    shape.shadow.inherit = False
    if rounded:
        try:
            shape.adjustments[0] = radius
        except Exception:
            pass
    return shape


def add_oval(slide, x, y, w, h, fill=BG_BLACK, line=None, line_width=Pt(0)):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = line_width
    shape.shadow.inherit = False
    return shape


def add_up_arrow(slide, x, y, w, h, fill=PALO_GREEN):
    shape = slide.shapes.add_shape(MSO_SHAPE.UP_ARROW, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_text(slide, x, y, w, h, runs, *, default_size=12, default_color=TEXT_WHITE,
             default_bold=False, default_font="Calibri", align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, line_spacing=1.15):
    """runs may be: a string, a list of strings, or a list of {text, bold, size, color, font} dicts."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = Pt(2)
    tf.margin_right = Pt(2)
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)
    tf.word_wrap = True
    tf.vertical_anchor = anchor

    if isinstance(runs, str):
        runs = [runs]

    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    for i, item in enumerate(runs):
        run = p.add_run()
        if isinstance(item, str):
            run.text = item
            run.font.name = default_font
            run.font.size = Pt(default_size)
            run.font.bold = default_bold
            run.font.color.rgb = default_color
        else:
            run.text = item.get("text", "")
            run.font.name = item.get("font", default_font)
            run.font.size = Pt(item.get("size", default_size))
            run.font.bold = item.get("bold", default_bold)
            run.font.color.rgb = item.get("color", default_color)
    return tb


def add_multiline(slide, x, y, w, h, paragraphs, *, size=10, color=TEXT_DIM, line_spacing=1.25,
                  font="Calibri", align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = Pt(2); tf.margin_right = Pt(2); tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
    tf.word_wrap = True
    for i, item in enumerate(paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        p.text = item if isinstance(item, str) else item.get("text", "")
        for run in p.runs:
            run.font.name = font
            run.font.size = Pt(size)
            run.font.color.rgb = color
            if isinstance(item, dict) and item.get("bold"):
                run.font.bold = True
    return tb


# --- Slide builder ----------------------------------------------------------
def build_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_bg(slide)

    # Layout regions — main area (left ~10.3"), right sidebar (~2.7" wide)
    MAIN_X = Inches(0.4)
    MAIN_W = Inches(10.2)
    SIDE_X = Inches(10.75)
    SIDE_W = Inches(2.3)

    # ----- TITLE BAR (top-left, two-tone: green + white) --------------------
    add_text(
        slide, MAIN_X, Inches(0.25), Inches(12), Inches(0.7),
        [
            {"text": "Cortex XDR: ", "bold": True, "size": 28, "color": PALO_GREEN},
            {"text": "Incident Response — One Case, Five Surfaces", "bold": True, "size": 28, "color": TEXT_WHITE},
        ],
        line_spacing=1.0,
    )

    # ----- CLOUD BACKDROP (subtle gray ellipses behind the card row) --------
    # Three overlapping clouds give the same airy backdrop the source slide uses
    add_oval(slide, Inches(0.4), Inches(1.9), Inches(4.5), Inches(2.5), fill=CLOUD_DARK)
    add_oval(slide, Inches(3.2), Inches(1.7), Inches(4.5), Inches(2.7), fill=CLOUD_GRAY)
    add_oval(slide, Inches(6.0), Inches(1.9), Inches(4.5), Inches(2.5), fill=CLOUD_DARK)

    # ----- CENTRAL WHITE BADGE ("THE CASE") ---------------------------------
    badge_w = Inches(2.0)
    badge_h = Inches(2.0)
    badge_cx = MAIN_X + MAIN_W / 2  # center of main area
    badge_x = badge_cx - badge_w / 2
    badge_y = Inches(0.95)
    # Green outer ring (slightly larger oval underneath)
    add_oval(slide, badge_x - Inches(0.05), badge_y - Inches(0.05),
             badge_w + Inches(0.1), badge_h + Inches(0.1),
             fill=PALO_GREEN)
    # White face
    add_oval(slide, badge_x, badge_y, badge_w, badge_h, fill=TEXT_WHITE)
    # Badge label
    add_text(
        slide, badge_x, badge_y + Inches(0.45), badge_w, Inches(0.5),
        [{"text": "THE", "bold": True, "size": 14, "color": GREEN_DIM}],
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide, badge_x, badge_y + Inches(0.8), badge_w, Inches(0.6),
        [{"text": "CASE", "bold": True, "size": 28, "color": BG_BLACK}],
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide, badge_x, badge_y + Inches(1.35), badge_w, Inches(0.5),
        [{"text": "in Cortex XDR", "bold": False, "size": 9, "color": GREEN_DIM}],
        align=PP_ALIGN.CENTER,
    )

    # ----- 5 CAPABILITY CARDS (horizontal row) ------------------------------
    cards = [
        ("☰", "Issues & Insights",
         ["Sortable table of all issues", "Low-severity insights surfaced", "Drill to detail"]),
        ("⧖", "Case Timeline",
         ["Issues added via causality", "Analyst action log", "Chronological view"]),
        ("⋈", "Causality Chain",
         ["Process tree + forensics highlights", "Event details + observables", "Live Terminal on endpoint", "Response actions from chain"]),
        ("⛔", "Response & Containment",
         ["Block file hash globally", "Isolate endpoint", "Quarantine artifact"]),
        ("⁂", "Collaboration & Automation",
         ["Case + Issue war rooms", "Playbooks for response", "Multi-analyst handoffs"]),
    ]

    card_count = len(cards)
    card_w = Inches(1.85)
    card_h = Inches(2.0)
    gap = Inches(0.05)
    cards_total_w = card_count * card_w + (card_count - 1) * gap
    cards_x_start = MAIN_X + (MAIN_W - cards_total_w) / 2
    cards_y = Inches(3.05)

    for i, (icon, title, bullets) in enumerate(cards):
        x = cards_x_start + i * (card_w + gap)
        # Card body
        add_rect(slide, x, cards_y, card_w, card_h, fill=CARD_DARK, rounded=True, radius=0.06,
                 line=CARD_OUTLINE, line_width=Pt(0.5))
        # Icon circle (overlaps top of card)
        icon_d = Inches(0.85)
        icon_cx = x + card_w / 2
        icon_x = icon_cx - icon_d / 2
        icon_y = cards_y - icon_d / 2
        # Green ring
        add_oval(slide, icon_x - Inches(0.04), icon_y - Inches(0.04),
                 icon_d + Inches(0.08), icon_d + Inches(0.08),
                 fill=PALO_GREEN)
        # Dark face
        add_oval(slide, icon_x, icon_y, icon_d, icon_d, fill=CARD_DARK)
        # Icon glyph
        add_text(
            slide, icon_x, icon_y, icon_d, icon_d,
            [{"text": icon, "bold": True, "size": 22, "color": PALO_GREEN}],
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
        )

        # Card title (extra height to absorb 2-line titles cleanly)
        add_text(
            slide, x + Inches(0.1), cards_y + Inches(0.48), card_w - Inches(0.2), Inches(0.6),
            [{"text": title, "bold": True, "size": 12, "color": TEXT_WHITE}],
            align=PP_ALIGN.CENTER,
        )

        # Underline accent (pushed below the title)
        underline_w = Inches(0.4)
        add_rect(slide, x + card_w / 2 - underline_w / 2, cards_y + Inches(1.12),
                 underline_w, Pt(1.5), fill=PALO_GREEN)

        # Bullets
        bullet_paras = [f"· {b}" for b in bullets]
        add_multiline(
            slide, x + Inches(0.12), cards_y + Inches(1.22),
            card_w - Inches(0.24), card_h - Inches(1.25),
            bullet_paras, size=8.5, color=TEXT_DIM, line_spacing=1.25,
        )

    # ----- UP-ARROWS from bottom strip to cards -----------------------------
    arrow_w = Inches(0.4)
    arrow_h = Inches(0.32)
    arrow_y = cards_y + card_h + Inches(0.08)
    for i in range(card_count):
        x = cards_x_start + i * (card_w + gap) + card_w / 2 - arrow_w / 2
        add_up_arrow(slide, x, arrow_y, arrow_w, arrow_h, fill=PALO_GREEN)

    # ----- BOTTOM STRIP — pivots / inputs that feed the case ---------------
    strip_x = MAIN_X
    strip_w = MAIN_W
    strip_y = cards_y + card_h + Inches(0.55)
    strip_h = Inches(1.25)
    add_rect(slide, strip_x, strip_y, strip_w, strip_h, fill=BG_BLACK,
             line=PALO_GREEN, line_width=Pt(1.5), rounded=True, radius=0.04)

    # (no strip label — the arrows + green-outlined strip already communicate
    # the pivot relationship; an explicit label collided with the up-arrows)

    # 5 pivot items inside the strip (same x-centers as the cards above)
    pivots = [
        ("▤", "Users & Assets",          "Entity pivot from case"),
        ("⌗", "File Hash View",          "Artifact-centric drill"),
        ("☰", "XQL Search",              "Hunt sibling activity"),
        ("⌖", "Forensics Module",        "Deep-triage on demand"),
        ("⌬", "MITRE Tactics + Tech.",   "ATT&CK map of issues"),
    ]
    item_y = strip_y + Inches(0.1)
    for i, (icon, name, sub) in enumerate(pivots):
        x = cards_x_start + i * (card_w + gap)
        # Icon
        add_text(
            slide, x, item_y, card_w, Inches(0.42),
            [{"text": icon, "bold": True, "size": 19, "color": PALO_GREEN}],
            align=PP_ALIGN.CENTER,
        )
        # Name
        add_text(
            slide, x, item_y + Inches(0.46), card_w, Inches(0.3),
            [{"text": name, "bold": True, "size": 10, "color": TEXT_WHITE}],
            align=PP_ALIGN.CENTER,
        )
        # Sub
        add_text(
            slide, x, item_y + Inches(0.74), card_w, Inches(0.3),
            [{"text": sub, "bold": False, "size": 8, "color": TEXT_MUTED}],
            align=PP_ALIGN.CENTER,
        )

    # ----- RIGHT SIDEBAR (3 value-prop blocks with divider lines) ----------
    # Top divider
    add_rect(slide, SIDE_X, Inches(1.3), SIDE_W, Pt(1), fill=GREEN_DIM)
    blocks = [
        ("Find the Root Cause", "in One Screen"),
        ("Act Without Leaving", "the Console"),
        ("Collaborate +", "Automate Response"),
    ]
    block_h = Inches(1.55)
    block_y = Inches(1.4)
    for i, (line1, line2) in enumerate(blocks):
        y = block_y + i * (block_h + Inches(0.25))
        add_text(
            slide, SIDE_X, y, SIDE_W, Inches(0.55),
            [{"text": line1, "bold": True, "size": 17, "color": TEXT_WHITE}],
            align=PP_ALIGN.LEFT,
        )
        add_text(
            slide, SIDE_X, y + Inches(0.55), SIDE_W, Inches(0.55),
            [{"text": line2, "bold": True, "size": 17, "color": TEXT_WHITE}],
            align=PP_ALIGN.LEFT,
        )
        # Divider after each except last
        if i < len(blocks) - 1:
            add_rect(slide, SIDE_X, y + block_h, SIDE_W * 0.8, Pt(1), fill=GREEN_DIM)

    # Bottom note (small, placed below the strip but still on-slide)
    add_text(
        slide, MAIN_X, Inches(7.18), Inches(12.5), Inches(0.25),
        [{"text": "Cortex XDR IR Workshop  ·  Investigation surfaces analysts use to work an incident case end-to-end",
          "size": 8.5, "color": TEXT_MUTED}],
    )


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    build_slide(prs)

    out = "/Users/ayman/Documents/Coding/phantom/docs/workshops/cortex-xdr-ir-surfaces.pptx"
    prs.save(out)
    print(f"Wrote {out}")


if __name__ == "__main__":
    main()
