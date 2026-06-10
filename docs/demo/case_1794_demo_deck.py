"""
Generates `case_1794_demo_deck.pptx` — the Case 1794 XDR investigation
demo deck. Run: `python3 case_1794_demo_deck.py`. Output sits next to
this script. Re-run after editing to regenerate; the .pptx is not
hand-edited so the script is the source of truth.

Style:
  - Dark slide background (security-tool aesthetic + projector-friendly)
  - White body text; red for "attack indicators / wow moments"; amber
    for "key pivot / decision moment"
  - Monospace blocks for XQL + commands
  - 16:9 widescreen, 13.333" × 7.5" (PowerPoint default)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pathlib import Path

# ─── Palette ────────────────────────────────────────────────────────
BG          = RGBColor(0x0F, 0x14, 0x1A)   # near-black, slight blue
PANEL       = RGBColor(0x1A, 0x22, 0x2C)   # lighter panel
WHITE       = RGBColor(0xEC, 0xEF, 0xF4)
MUTED       = RGBColor(0x9A, 0xA5, 0xB1)
RED         = RGBColor(0xFF, 0x4D, 0x5E)   # attack / wow
AMBER       = RGBColor(0xFF, 0xB4, 0x4D)   # pivot / decision
GREEN       = RGBColor(0x6E, 0xCB, 0x8C)   # success / good signal
BLUE        = RGBColor(0x4F, 0xA8, 0xFF)   # info / linkish
PURPLE      = RGBColor(0xB4, 0x8A, 0xFF)   # XQL / advanced

FONT_BODY   = "Helvetica Neue"
FONT_MONO   = "Menlo"

# ─── Helpers ────────────────────────────────────────────────────────

def set_slide_bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, left, top, width, height, text, *,
             size=18, bold=False, color=WHITE, font=FONT_BODY,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    """Add a single styled text box. Multi-line if `text` has `\n`."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)

    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return tb


def add_runs(slide, left, top, width, height, runs, *,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=None):
    """Add a textbox with multiple styled inline runs.
    `runs` is a list of either strings (newline = new paragraph) or
    dicts {text, size, bold, color, font, newline}.
    """
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)

    # Start the first paragraph
    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing:
        p.line_spacing = line_spacing
    first = True
    for r in runs:
        if isinstance(r, str):
            # Treat \n in raw strings as paragraph break for ergonomics
            parts = r.split("\n")
            for j, part in enumerate(parts):
                if j > 0:
                    p = tf.add_paragraph()
                    p.alignment = align
                    if line_spacing:
                        p.line_spacing = line_spacing
                if part:
                    run = p.add_run()
                    run.text = part
                    run.font.name = FONT_BODY
                    run.font.size = Pt(16)
                    run.font.color.rgb = WHITE
        else:
            if r.get("newline"):
                p = tf.add_paragraph()
                p.alignment = align
                if line_spacing:
                    p.line_spacing = line_spacing
                continue
            run = p.add_run()
            run.text = r.get("text", "")
            run.font.name = r.get("font", FONT_BODY)
            run.font.size = Pt(r.get("size", 16))
            run.font.bold = r.get("bold", False)
            run.font.color.rgb = r.get("color", WHITE)
    return tb


def add_panel(slide, left, top, width, height, *,
              fill=PANEL, line=None, line_color=MUTED):
    """Background rectangle for grouping. No text."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line)
    # Mild rounding via adjustments — keep default
    return shape


def add_accent_bar(slide, top, color=RED, width=Inches(0.08), height=None):
    """Left-edge vertical accent bar. Defaults to full slide height."""
    if height is None:
        height = Inches(7.5) - top
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.4), top, width, height,
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    return bar


def slide_title(slide, eyebrow, title, *,
                eyebrow_color=AMBER, title_color=WHITE):
    """Standard slide header: small uppercase eyebrow + big title."""
    add_text(slide, Inches(0.7), Inches(0.4), Inches(12), Inches(0.35),
             eyebrow.upper(), size=12, bold=True, color=eyebrow_color)
    add_text(slide, Inches(0.7), Inches(0.75), Inches(12), Inches(0.7),
             title, size=28, bold=True, color=title_color)


def code_block(slide, left, top, width, height, lines, *,
               size=12, line_color=MUTED):
    """Monospace code block with subtle background panel."""
    add_panel(slide, left, top, width, height, fill=RGBColor(0x10, 0x18, 0x22),
              line=0.5, line_color=line_color)
    text_box = add_text(slide, left + Inches(0.12), top + Inches(0.1),
                        width - Inches(0.24), height - Inches(0.2),
                        "\n".join(lines), size=size, color=GREEN,
                        font=FONT_MONO)
    return text_box


# ─── Build the deck ─────────────────────────────────────────────────

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]  # blank layout


# ─── Slide 1: Title ─────────────────────────────────────────────────
s = prs.slides.add_slide(blank); set_slide_bg(s)
add_accent_bar(s, Inches(0))
add_text(s, Inches(0.7), Inches(2.2), Inches(12), Inches(0.5),
         "DEMO WALKTHROUGH", size=14, bold=True, color=AMBER)
add_text(s, Inches(0.7), Inches(2.75), Inches(12), Inches(1.2),
         "Case 1794", size=64, bold=True, color=WHITE)
add_text(s, Inches(0.7), Inches(3.95), Inches(12), Inches(0.7),
         "Phishing → Ransomware Kill Chain", size=32, color=WHITE)
add_text(s, Inches(0.7), Inches(4.7), Inches(12), Inches(0.5),
         "A cohesive XDR investigation flow — not console hopping",
         size=18, color=MUTED)
add_text(s, Inches(0.7), Inches(6.6), Inches(12), Inches(0.4),
         "22 steps · 7 Acts · 2 endpoints · 1 incident", size=14, color=MUTED)


# ─── Slide 2: The framing ──────────────────────────────────────────
s = prs.slides.add_slide(blank); set_slide_bg(s)
add_accent_bar(s, Inches(0), color=AMBER)
slide_title(s, "Set the scene", "Before clicking anything", eyebrow_color=AMBER)

add_panel(s, Inches(0.7), Inches(1.7), Inches(11.9), Inches(3.3),
          fill=PANEL)
add_text(s, Inches(1.0), Inches(1.95), Inches(11.3), Inches(3.0),
         ('"Case 1794 starts with a single phishing email and ends with a '
          'wiped Security event log on two hosts.\n\n'
          '22 steps, 7 acts, 2 compromised endpoints. What I\'m going to '
          'show you is how XDR turns that into ONE incident — not 22 '
          'disconnected alerts — and how an analyst pivots through it in '
          'about 8 minutes to get from \'what happened\' to \'who, what, '
          'how, and what\'s the blast radius.\'\n\n'
          'We\'ll move through it the way the attacker did, top to bottom."'),
         size=18, color=WHITE)

add_text(s, Inches(0.7), Inches(5.3), Inches(12), Inches(0.4),
         "WHY THIS FRAMING MATTERS",
         size=12, bold=True, color=AMBER)
add_text(s, Inches(0.7), Inches(5.7), Inches(12), Inches(1.5),
         ("The audience is now primed to evaluate \"did XDR connect the "
          "dots?\" — not \"did each rule fire?\""), size=18, color=WHITE)


# ─── Slide 3: Attack chain (reference to operator's image) ─────────
s = prs.slides.add_slide(blank); set_slide_bg(s)
add_accent_bar(s, Inches(0), color=BLUE)
slide_title(s, "Reference", "The attack chain we'll walk", eyebrow_color=BLUE)

# Acts table
acts = [
    ("Act 1 · Initial Foothold",   "steps 1–5",  "Phishing → cmd → PowerShell → bitsadmin → mars.exe beacons home"),
    ("Act 2 · Learning the Ground","steps 6–9",  "Discovery burst + cached creds + LSASS dump"),
    ("Act 3 · Getting Root",       "steps 10–12","UAC bypass · Defender silenced · certutil LOLBin"),
    ("Act 4 · Staying Put",        "steps 13–15","Local user · Run key · Scheduled task"),
    ("Act 5 · Lateral Movement",   "step 16",    "Stolen creds → SMB + WMI + WinRM → victim2 alive"),
    ("Act 6 · Harvest + Exfil",    "steps 17–20","Collect sensitive files · archive · DNS tunnel · HTTP POST"),
    ("Act 7 · Impact + Cleanup",   "steps 21–22","Defacement · clear Security event log (1102)"),
]
top = Inches(1.7)
row_h = Inches(0.62)
for i, (act, steps, desc) in enumerate(acts):
    y = top + row_h * i
    # Number badge
    badge = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.7), y + Inches(0.1),
                               Inches(0.42), Inches(0.42))
    badge.fill.solid(); badge.fill.fore_color.rgb = RED if i in (0, 2, 4, 6) else AMBER
    badge.line.fill.background()
    tf = badge.text_frame
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    r = tf.paragraphs[0].add_run()
    r.text = str(i + 1)
    r.font.name = FONT_BODY; r.font.size = Pt(16); r.font.bold = True
    r.font.color.rgb = BG
    # Act name
    add_text(s, Inches(1.25), y + Inches(0.05), Inches(3.5), Inches(0.4),
             act, size=15, bold=True, color=WHITE)
    # Steps tag
    add_text(s, Inches(4.8), y + Inches(0.05), Inches(1.5), Inches(0.4),
             steps, size=12, color=MUTED, font=FONT_MONO)
    # Description
    add_text(s, Inches(6.4), y + Inches(0.05), Inches(6.6), Inches(0.4),
             desc, size=13, color=WHITE)


# ─── Slide 4: Demo time budget ─────────────────────────────────────
s = prs.slides.add_slide(blank); set_slide_bg(s)
add_accent_bar(s, Inches(0), color=AMBER)
slide_title(s, "Pacing", "Time budget — ~14 min total", eyebrow_color=AMBER)

budget = [
    ("Framing", "1 min", "Set the story arc"),
    ("Act 1 — Foothold (Causality)", "1.5 min", ""),
    ("Act 2 — Discovery + first XQL", "1.5 min", ""),
    ("Act 3 — Defender silenced", "2 min", "🔥 wow moment 1 · slow down"),
    ("Act 4 — Persistence", "1 min", ""),
    ("Act 5 — Lateral pivot", "2 min", "🔥 wow moment 2 · slow down"),
    ("Act 6 — Exfil + XQL hunts", "1.5 min", ""),
    ("Act 7 — Security 1102 backstop", "1.5 min", ""),
    ("Response demo", "2 min", "Isolate · Live Terminal · Playbook"),
    ("Q&A buffer", "~1 min", ""),
]
top = Inches(1.7)
row_h = Inches(0.48)
for i, (item, dur, note) in enumerate(budget):
    y = top + row_h * i
    add_text(s, Inches(0.9), y, Inches(5.5), Inches(0.4),
             item, size=15, bold=True, color=WHITE)
    add_text(s, Inches(6.6), y, Inches(1.5), Inches(0.4),
             dur, size=15, color=AMBER, font=FONT_MONO)
    if note:
        add_text(s, Inches(8.4), y, Inches(4.5), Inches(0.4),
                 note, size=13,
                 color=RED if "wow" in note else MUTED)


# ─── Slide 5: Pre-demo prep ────────────────────────────────────────
s = prs.slides.add_slide(blank); set_slide_bg(s)
add_accent_bar(s, Inches(0), color=BLUE)
slide_title(s, "Pre-flight (5 min before)", "Set up tabs · save queries · verify state", eyebrow_color=BLUE)

prep = [
    ("Open Case 1794",
     "Incident Response → Incidents → search “1794”",
     "Lands you on the Incident Card — the canonical entry point"),
    ("Verify XDR agents online",
     "Endpoints → Endpoint Administration",
     "xdragent (victim1) + xdragent2 (victim2)"),
    ("Save XQL queries to Library",
     "Investigation → Query Builder → Saved Queries",
     "Have queries #1–#7 (later slides) saved + named"),
    ("Pre-bookmark 3 pivot tabs",
     "Asset (victim1) · User (phantomlab) · Hash (mars.exe)",
     "Tab muscle-memory > “where was that link again”"),
]
top = Inches(1.8)
row_h = Inches(1.2)
for i, (action, where, why) in enumerate(prep):
    y = top + row_h * i
    add_panel(s, Inches(0.7), y, Inches(11.9), Inches(1.0), fill=PANEL)
    add_text(s, Inches(0.95), y + Inches(0.15), Inches(4.2), Inches(0.4),
             action, size=16, bold=True, color=WHITE)
    add_text(s, Inches(0.95), y + Inches(0.55), Inches(11.4), Inches(0.4),
             where, size=12, color=AMBER, font=FONT_MONO)
    add_text(s, Inches(5.4), y + Inches(0.15), Inches(7.0), Inches(0.4),
             why, size=13, color=MUTED)


# ─── Slide 6: Act 1 title ──────────────────────────────────────────
s = prs.slides.add_slide(blank); set_slide_bg(s)
add_accent_bar(s, Inches(0), color=RED)
slide_title(s, "Act 1 · Initial Foothold (steps 1–5)",
            "Alert Grouping Graph → Causality View root", eyebrow_color=RED)

add_text(s, Inches(0.7), Inches(1.8), Inches(12), Inches(0.4),
         "CONSOLE ACTION", size=12, bold=True, color=AMBER)
add_text(s, Inches(0.7), Inches(2.15), Inches(12), Inches(0.7),
         "Incident Card → Alert Grouping Graph (top-right)",
         size=20, bold=True, color=WHITE, font=FONT_MONO)

add_text(s, Inches(0.7), Inches(3.0), Inches(12), Inches(0.4),
         "WHAT TO SAY", size=12, bold=True, color=AMBER)
add_text(s, Inches(0.7), Inches(3.4), Inches(12), Inches(3.6),
         ("\"This is ONE incident, but it contains [N] alerts. The Alert "
          "Grouping Graph shows how XDR's Causality Analysis Engine bound "
          "them together.\n\n"
          "Look at the root node: cmd.exe spawned by Word's macro. That's "
          "the Causality Group Owner — the process the engine identified "
          "as responsible for everything downstream.\n\n"
          "Every alert on this graph traces back here.\""),
         size=18, color=WHITE)


# ─── Slide 7: Act 1 causality walk ─────────────────────────────────
s = prs.slides.add_slide(blank); set_slide_bg(s)
add_accent_bar(s, Inches(0), color=RED)
slide_title(s, "Act 1 — Causality View walk", "Root → child → child", eyebrow_color=RED)

walk = [
    ("1", "WINWORD.EXE",
     "Root node — user opens the attached invoice. Hover the Insights badge: \"macro execution from email-staged document.\""),
    ("2", "cmd.exe",
     "Word macro spawn. Right-click → Add to Timeline for later timestamp reference."),
    ("3", "powershell.exe (encoded)",
     "Encoded command argument. Hover the node — XDR auto-deobfuscates in the side panel."),
    ("4", "bitsadmin.exe",
     "LOLBin used to fetch nars.exe. Right-click → Hash View on mars.exe hash."),
    ("5", "mars.exe beacons home",
     '⚡ "MOMENT OF INITIAL COMPROMISE" anchor — note timestamp.'),
]
top = Inches(1.7)
row_h = Inches(0.95)
for i, (n, name, desc) in enumerate(walk):
    y = top + row_h * i
    badge = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.7), y + Inches(0.18),
                               Inches(0.5), Inches(0.5))
    badge.fill.solid(); badge.fill.fore_color.rgb = RED
    badge.line.fill.background()
    tf = badge.text_frame
    tf.margin_left = tf.margin_right = Emu(0); tf.margin_top = tf.margin_bottom = Emu(0)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    r = tf.paragraphs[0].add_run(); r.text = n
    r.font.name = FONT_BODY; r.font.size = Pt(20); r.font.bold = True
    r.font.color.rgb = BG
    add_text(s, Inches(1.35), y + Inches(0.1), Inches(4), Inches(0.45),
             name, size=18, bold=True, color=WHITE, font=FONT_MONO)
    add_text(s, Inches(1.35), y + Inches(0.5), Inches(11.5), Inches(0.5),
             desc, size=13, color=MUTED)

add_text(s, Inches(0.7), Inches(6.85), Inches(12), Inches(0.5),
         '💡  "Notice I haven\'t typed any XQL yet. Everything for the foothold story is in Causality View."',
         size=14, color=AMBER)


# ─── Slide 8: Act 2 — Discovery + Timeline + first XQL ─────────────
s = prs.slides.add_slide(blank); set_slide_bg(s)
add_accent_bar(s, Inches(0), color=PURPLE)
slide_title(s, "Act 2 · Learning the Ground (steps 6–9)",
            "Timeline View → first XQL", eyebrow_color=PURPLE)

add_text(s, Inches(0.7), Inches(1.7), Inches(6.0), Inches(0.4),
         "CONSOLE ACTION", size=12, bold=True, color=AMBER)
add_text(s, Inches(0.7), Inches(2.05), Inches(6.0), Inches(0.6),
         "Alert → Timeline tab",
         size=18, bold=True, color=WHITE, font=FONT_MONO)
add_text(s, Inches(0.7), Inches(2.85), Inches(6.0), Inches(3.4),
         ("\"Act 2 is the discovery burst. Six commands in 90 seconds — "
          "whoami, net share, net user, systeminfo, arp, netstat. "
          "Individually benign. The PATTERN — six discovery commands in "
          "under two minutes from the same PowerShell session — is what "
          "XDR's behavioral analytics caught.\""),
         size=14, color=WHITE)
add_text(s, Inches(0.7), Inches(5.5), Inches(6.0), Inches(0.4),
         "Open Forensics Highlights panel:", size=13, color=AMBER)
add_text(s, Inches(0.7), Inches(5.85), Inches(6.0), Inches(1.5),
         "• Account Discovery badge\n• System Information Discovery badge\n• Network Share Discovery badge\n  → each links to its MITRE technique",
         size=13, color=MUTED)

# XQL on right
add_text(s, Inches(7.0), Inches(1.7), Inches(6.0), Inches(0.4),
         "XQL #1 — DISCOVERY BURST", size=12, bold=True, color=PURPLE)
code_block(s, Inches(7.0), Inches(2.1), Inches(5.8), Inches(4.5), [
    "dataset = xdr_data",
    "| filter event_type = ENUM.PROCESS",
    "    and event_sub_type = ENUM.PROCESS_START",
    "    and agent_hostname = \"xdragent\"",
    "    and actor_process_image_name in",
    "        (\"powershell.exe\", \"cmd.exe\")",
    "    and action_process_image_name in (",
    "        \"whoami.exe\", \"net.exe\", \"net1.exe\",",
    "        \"systeminfo.exe\", \"arp.exe\",",
    "        \"netstat.exe\", \"route.exe\",",
    "        \"cmdkey.exe\"",
    "    )",
    "| fields _time,",
    "    action_process_image_name,",
    "    action_process_image_command_line,",
    "    actor_process_image_name",
    "| sort _time asc",
], size=11)
add_text(s, Inches(7.0), Inches(6.75), Inches(6), Inches(0.4),
         "Pivot to a hunt: remove agent_hostname filter → search all hosts.",
         size=12, color=AMBER)


# ─── Slide 9: Credential dumping + LSASS ──────────────────────────
s = prs.slides.add_slide(blank); set_slide_bg(s)
add_accent_bar(s, Inches(0), color=RED)
slide_title(s, "Act 2 (cont.) · Credential dumping (steps 7–9)",
            "The marquee EDR detection event", eyebrow_color=RED)

add_text(s, Inches(0.7), Inches(1.8), Inches(12), Inches(0.4),
         "FOLLOW THE CAUSALITY CHAIN TO:", size=12, bold=True, color=AMBER)
chain = [
    ("cmdkey.exe", "Cached credentials access (T1003.005)"),
    ("rundll32.exe comsvcs.dll", "LSASS memory dump via MiniDump trick — T1003.001"),
    ("PowerKatz invocation", "In-memory credential extraction"),
]
for i, (cmd, desc) in enumerate(chain):
    y = Inches(2.4) + Inches(0.7) * i
    add_text(s, Inches(0.9), y, Inches(4.5), Inches(0.5),
             cmd, size=18, bold=True, color=RED, font=FONT_MONO)
    add_text(s, Inches(5.5), y + Inches(0.05), Inches(7.5), Inches(0.5),
             desc, size=15, color=WHITE)

add_panel(s, Inches(0.7), Inches(5.0), Inches(11.9), Inches(2.0), fill=PANEL)
add_text(s, Inches(1.0), Inches(5.2), Inches(11.4), Inches(0.5),
         "👉 For LSASS specifically:", size=15, bold=True, color=AMBER)
add_text(s, Inches(1.0), Inches(5.7), Inches(11.4), Inches(1.3),
         ("Click the LSASS dump process node → All Events tab → show the "
          "process_access_image_load event where rundll32 opened a handle "
          "to lsass.exe (XDR's equivalent of Sysmon EID 10). "
          "This is the marquee EDR detection signal."),
         size=14, color=WHITE)


# ─── Slide 10: Act 3 — DEFENDER SILENCED (wow moment 1) ────────────
s = prs.slides.add_slide(blank); set_slide_bg(s)
add_accent_bar(s, Inches(0), color=RED)
add_text(s, Inches(0.7), Inches(0.4), Inches(12), Inches(0.4),
         "ACT 3 · GETTING ROOT (STEPS 10–12)", size=12, bold=True, color=RED)
add_text(s, Inches(0.7), Inches(0.85), Inches(12), Inches(1.4),
         "DEFENDER SILENCED", size=56, bold=True, color=RED)
add_text(s, Inches(0.7), Inches(2.4), Inches(12), Inches(0.6),
         "...downstream telemetry degrades", size=24, color=AMBER)

add_panel(s, Inches(0.7), Inches(3.4), Inches(11.9), Inches(3.5), fill=PANEL)
add_text(s, Inches(1.0), Inches(3.6), Inches(11.3), Inches(3.2),
         ("\"Three things in 30 seconds: bypass UAC via Fodhelper, disable "
          "Microsoft Defender real-time monitoring, decode the next-stage "
          "payload with certutil.\n\n"
          "If we relied on Defender alone, we'd lose visibility right here.\n\n"
          "But Defender being silenced IS itself a high-severity "
          "Defender/Operational event 5001 — and XDR's behavioral telemetry "
          "continues regardless because it's collected by the Cortex XDR "
          "agent, not Defender.\""),
         size=17, color=WHITE)

add_text(s, Inches(0.7), Inches(7.05), Inches(12), Inches(0.4),
         "🔥  This is the demo's strongest moment — slow down here.",
         size=14, color=AMBER, bold=True)


# ─── Slide 11: Act 3 demonstrate parallel telemetry ────────────────
s = prs.slides.add_slide(blank); set_slide_bg(s)
add_accent_bar(s, Inches(0), color=RED)
slide_title(s, "Act 3 (cont.) · Show the parallel telemetry path",
            "Cortex agent ≠ Defender agent", eyebrow_color=RED)

steps = [
    ("Click", "“Defender real-time monitoring disabled” alert → Open Card"),
    ("Show", "Causality View — step 12 (certutil decode) STILL has full event detail AFTER Defender went off"),
    ("Forensics Highlights", "“Indicator removal — defense evasion” badge"),
]
top = Inches(1.8)
row_h = Inches(0.9)
for i, (verb, action) in enumerate(steps):
    y = top + row_h * i
    add_text(s, Inches(0.9), y, Inches(2.5), Inches(0.4),
             verb, size=15, bold=True, color=AMBER)
    add_text(s, Inches(3.5), y, Inches(9.3), Inches(0.7),
             action, size=15, color=WHITE)

add_text(s, Inches(0.7), Inches(4.7), Inches(12), Inches(0.4),
         "THEN PIVOT TO ASSET VIEW (HOST INSIGHTS)",
         size=13, bold=True, color=AMBER)
add_panel(s, Inches(0.7), Inches(5.15), Inches(11.9), Inches(1.7), fill=PANEL)
add_text(s, Inches(1.0), Inches(5.35), Inches(11.3), Inches(1.4),
         ("• Click Asset link on the alert → Host Insights for xdragent\n"
          "• Installed Software / Running Services — confirm Defender service state\n"
          "• Vulnerability Assessment panel — pre-existing CVEs that contextualize \"was this attractive prey?\""),
         size=14, color=WHITE)


# ─── Slide 12: Act 4 — Persistence ────────────────────────────────
s = prs.slides.add_slide(blank); set_slide_bg(s)
add_accent_bar(s, Inches(0), color=AMBER)
slide_title(s, "Act 4 · Staying Put (steps 13–15)",
            "Forensics Highlights · Live Terminal preview", eyebrow_color=AMBER)

add_text(s, Inches(0.7), Inches(1.7), Inches(12), Inches(0.4),
         "CONSOLE ACTION", size=12, bold=True, color=AMBER)
add_text(s, Inches(0.7), Inches(2.05), Inches(12), Inches(0.5),
         "Filter incident alerts → tactic: Persistence (left filter)",
         size=18, bold=True, color=WHITE, font=FONT_MONO)
add_text(s, Inches(0.7), Inches(2.75), Inches(12), Inches(0.7),
         ("\"Three persistence mechanisms in under a minute. Signal you "
          "want your SOC to internalize: multiple persistence mechanisms "
          "in burst → don't deal with them as isolated alerts.\""),
         size=15, color=WHITE)

# Three persistence cards
top = Inches(4.0)
items = [
    ("Local user create", "Security 4720", "Account: T1136.001_PowerShell"),
    ("Registry Run key", "Sysmon EID 13", "Value: PhantomUpdater → implant path"),
    ("Scheduled task", "Security 4698", "schtasks /create /sc onlogon /ru SYSTEM"),
]
for i, (title, evt, detail) in enumerate(items):
    x = Inches(0.7) + Inches(4.05) * i
    add_panel(s, x, top, Inches(3.9), Inches(2.3), fill=PANEL)
    add_text(s, x + Inches(0.25), top + Inches(0.2), Inches(3.6), Inches(0.5),
             title, size=16, bold=True, color=WHITE)
    add_text(s, x + Inches(0.25), top + Inches(0.75), Inches(3.6), Inches(0.4),
             evt, size=12, color=AMBER, font=FONT_MONO)
    add_text(s, x + Inches(0.25), top + Inches(1.2), Inches(3.6), Inches(1.0),
             detail, size=13, color=MUTED)


# ─── Slide 13: Live Terminal preview ──────────────────────────────
s = prs.slides.add_slide(blank); set_slide_bg(s)
add_accent_bar(s, Inches(0), color=AMBER)
slide_title(s, "Act 4 · Preview Live Terminal (don't run yet)",
            "Initiate Response Action → Live Terminal", eyebrow_color=AMBER)

add_panel(s, Inches(0.7), Inches(1.8), Inches(11.9), Inches(4.5), fill=PANEL)
add_text(s, Inches(1.0), Inches(2.0), Inches(11.3), Inches(4.3),
         ("Click the Initiate Response Action button on victim1 → Live Terminal.\n\n"
          "Show the welcome screen — but DO NOT run anything yet.\n\n"
          "\"From here I could run PowerShell, kill the scheduled task, delete "
          "the Run key, dump memory. I can do all of that without RDP-ing in "
          "or coordinating with the desktop team.\n\n"
          "We'll come back to this in the response phase at the end.\""),
         size=16, color=WHITE)

add_text(s, Inches(0.7), Inches(6.5), Inches(12), Inches(0.5),
         "📌  Foreshadowing builds anticipation. Operators love this feature — tease it.",
         size=14, color=AMBER)


# ─── Slide 14: Act 5 — Lateral movement title ─────────────────────
s = prs.slides.add_slide(blank); set_slide_bg(s)
add_accent_bar(s, Inches(0), color=RED)
add_text(s, Inches(0.7), Inches(0.4), Inches(12), Inches(0.4),
         "ACT 5 · LATERAL MOVEMENT (STEP 16)", size=12, bold=True, color=RED)
add_text(s, Inches(0.7), Inches(0.85), Inches(12), Inches(1.4),
         "THE CRITICAL PIVOT", size=56, bold=True, color=RED)
add_text(s, Inches(0.7), Inches(2.4), Inches(12), Inches(0.6),
         "blast radius doubles", size=24, color=AMBER)

add_panel(s, Inches(0.7), Inches(3.4), Inches(11.9), Inches(3.4), fill=PANEL)
add_text(s, Inches(1.0), Inches(3.6), Inches(11.3), Inches(3.1),
         ("Step 16: attacker on victim1 with elevated creds reaches across "
          "to victim2 via SMB + WMI + WinRM.\n\n"
          "Console: Open Network Causality View on the lateral alert.\n\n"
          "\"This is the part XDR does that legacy AV simply cannot. The "
          "destination host has its own agent feeding XDR — we'll see the "
          "remote process creation from victim2's perspective, proving "
          "the lateral worked.\""),
         size=17, color=WHITE)

add_text(s, Inches(0.7), Inches(7.0), Inches(12), Inches(0.4),
         "🔥  Wow moment 2 — slow down, give it the dramatic beat.",
         size=14, color=AMBER, bold=True)


# ─── Slide 15: Network Causality View walk ─────────────────────────
s = prs.slides.add_slide(blank); set_slide_bg(s)
add_accent_bar(s, Inches(0), color=RED)
slide_title(s, "Act 5 · Network Causality View walk",
            "victim1 → victim2 cross-host evidence", eyebrow_color=RED)

# Source node
add_panel(s, Inches(0.7), Inches(1.8), Inches(3.5), Inches(2.2),
          fill=PANEL, line=1, line_color=RED)
add_text(s, Inches(0.85), Inches(1.95), Inches(3.2), Inches(0.4),
         "SOURCE", size=11, bold=True, color=RED)
add_text(s, Inches(0.85), Inches(2.35), Inches(3.2), Inches(0.5),
         "xdragent", size=22, bold=True, color=WHITE, font=FONT_MONO)
add_text(s, Inches(0.85), Inches(2.85), Inches(3.2), Inches(0.4),
         "10.10.0.14", size=14, color=MUTED, font=FONT_MONO)
add_text(s, Inches(0.85), Inches(3.3), Inches(3.2), Inches(0.7),
         "PowerShell session\norigin of lateral", size=13, color=WHITE)

# Arrow + protocols
add_text(s, Inches(4.5), Inches(2.4), Inches(4.3), Inches(0.5),
         "──────►", size=28, color=RED, align=PP_ALIGN.CENTER)
add_text(s, Inches(4.5), Inches(2.9), Inches(4.3), Inches(0.4),
         "TCP 445 (SMB)", size=14, color=AMBER, font=FONT_MONO, align=PP_ALIGN.CENTER)
add_text(s, Inches(4.5), Inches(3.3), Inches(4.3), Inches(0.4),
         "TCP 5985 (WinRM)", size=14, color=AMBER, font=FONT_MONO, align=PP_ALIGN.CENTER)
add_text(s, Inches(4.5), Inches(3.7), Inches(4.3), Inches(0.4),
         "user: phantomlab", size=14, color=AMBER, font=FONT_MONO, align=PP_ALIGN.CENTER)

# Dest node
add_panel(s, Inches(9.1), Inches(1.8), Inches(3.5), Inches(2.2),
          fill=PANEL, line=1, line_color=RED)
add_text(s, Inches(9.25), Inches(1.95), Inches(3.2), Inches(0.4),
         "DESTINATION", size=11, bold=True, color=RED)
add_text(s, Inches(9.25), Inches(2.35), Inches(3.2), Inches(0.5),
         "xdragent2", size=22, bold=True, color=WHITE, font=FONT_MONO)
add_text(s, Inches(9.25), Inches(2.85), Inches(3.2), Inches(0.4),
         "10.10.0.16", size=14, color=MUTED, font=FONT_MONO)
add_text(s, Inches(9.25), Inches(3.3), Inches(3.2), Inches(0.7),
         "wmiprvse.exe spawn\nmarker file write", size=13, color=WHITE)

# Pivot list
add_text(s, Inches(0.7), Inches(4.5), Inches(12), Inches(0.4),
         "NEXT PIVOTS FROM THIS VIEW:", size=12, bold=True, color=AMBER)
pivots = [
    "Destination node → Open Card → xdragent2's causality (wmiprvse → cmd → marker file write)",
    "Auth event → click User → phantomlab → User View (logon history across environment)",
    "Asset enrichment for accounts marked admin (this account shouldn't be lateral-active)",
]
for i, txt in enumerate(pivots):
    add_text(s, Inches(1.0), Inches(4.9) + Inches(0.45) * i, Inches(11.5), Inches(0.4),
             "• " + txt, size=14, color=WHITE)


# ─── Slide 16: Act 5 XQL ──────────────────────────────────────────
s = prs.slides.add_slide(blank); set_slide_bg(s)
add_accent_bar(s, Inches(0), color=PURPLE)
slide_title(s, "Act 5 · XQL #2 — cross-host auth proof",
            "Confirm what Network Causality showed visually", eyebrow_color=PURPLE)

code_block(s, Inches(0.7), Inches(1.8), Inches(11.9), Inches(2.6), [
    "dataset = authentication",
    "| filter dst_host_name = \"xdragent2\"",
    "    and src_host_name = \"xdragent\"",
    "    and action_username = \"phantomlab\"",
    "    and _time >= toepoch(<op_start>)",
    "| fields _time, authentication_protocol, action_username,",
    "    src_host_ip, dst_host_ip, action_outcome",
    "| sort _time desc",
], size=14)

add_panel(s, Inches(0.7), Inches(4.7), Inches(11.9), Inches(2.4), fill=PANEL)
add_text(s, Inches(1.0), Inches(4.85), Inches(11.3), Inches(2.2),
         ("Result will show: NTLM auth, succeeded, originating from "
          "10.10.0.14, targeting 10.10.0.16.\n\n"
          "\"Without cross-host correlation, my SOC sees TWO incidents — "
          "'suspicious process on victim1' and 'suspicious WMI activity on "
          "victim2.' With XDR, it's ONE incident with the lateral edge "
          "between them.\n\n"
          "That's the difference between an analyst noticing a real attack "
          "vs. closing two low-priority alerts as 'unrelated.'\""),
         size=14, color=WHITE)


# ─── Slide 17: Act 6 — Exfil + XQL hunts ───────────────────────────
s = prs.slides.add_slide(blank); set_slide_bg(s)
add_accent_bar(s, Inches(0), color=PURPLE)
slide_title(s, "Act 6 · Harvest + Exfil (steps 17–20)",
            "All Events view → XQL #3 + #4 → Retrieve File", eyebrow_color=PURPLE)

add_text(s, Inches(0.7), Inches(1.6), Inches(6.0), Inches(0.4),
         "XQL #3 — DNS TUNNELING", size=12, bold=True, color=PURPLE)
code_block(s, Inches(0.7), Inches(1.95), Inches(6.0), Inches(2.6), [
    "dataset = xdr_data",
    "| filter event_type = ENUM.DNS",
    "    and agent_hostname = \"xdragent2\"",
    "    and _time >= toepoch(<op_start>)",
    "| comp count_distinct(dns_query_name)",
    "    as unique_subdomains by",
    "    bin(_time, 1m),",
    "    regextract(dns_query_name,",
    "      \".*\\.(.*\\..*\\..*)$\") as parent",
    "| filter unique_subdomains >= 10",
    "| sort _time desc",
], size=11)

add_text(s, Inches(7.0), Inches(1.6), Inches(6.0), Inches(0.4),
         "XQL #4 — ARCHIVE IN TEMP", size=12, bold=True, color=PURPLE)
code_block(s, Inches(7.0), Inches(1.95), Inches(6.0), Inches(2.6), [
    "dataset = xdr_data",
    "| filter event_type = ENUM.FILE",
    "    and event_sub_type = ENUM.FILE_CREATE_NEW",
    "    and action_file_extension = \"zip\"",
    "    and action_file_path contains \"Temp\"",
    "    and actor_process_image_name",
    "      = \"powershell.exe\"",
    "| fields _time, agent_hostname,",
    "    action_file_path, action_file_size,",
    "    actor_process_command_line",
], size=11)

add_panel(s, Inches(0.7), Inches(4.8), Inches(11.9), Inches(2.2), fill=PANEL)
add_text(s, Inches(1.0), Inches(5.0), Inches(11.3), Inches(0.5),
         "FORENSIC ACTION — Retrieve File",
         size=15, bold=True, color=AMBER)
add_text(s, Inches(1.0), Inches(5.5), Inches(11.3), Inches(1.5),
         ("From any alert with the archive's path, click Retrieve File in "
          "the right-side action panel.\n\n"
          "\"In a real investigation, I'd pull this zip back to forensics "
          "for inspection — verify what was actually staged for exfil. "
          "Don't have to image the disk.\""),
         size=14, color=WHITE)


# ─── Slide 18: Act 7 — Security 1102 backstop ──────────────────────
s = prs.slides.add_slide(blank); set_slide_bg(s)
add_accent_bar(s, Inches(0), color=GREEN)
slide_title(s, "Act 7 · Impact + Cleanup (steps 21–22)",
            "Security event 1102 — the smoking gun nobody can hide", eyebrow_color=GREEN)

add_panel(s, Inches(0.7), Inches(1.7), Inches(11.9), Inches(3.0), fill=PANEL)
add_text(s, Inches(1.0), Inches(1.95), Inches(11.3), Inches(2.7),
         ("\"Security event 1102 fires unconditionally on log clear — "
          "Microsoft built it specifically to flag this anti-forensics "
          "pattern.\n\n"
          "Even if the attacker cleared the rest of the Security log, THIS "
          "event lands in XDR before it's gone, because the Cortex XDR "
          "agent forwards in real time.\""),
         size=18, color=WHITE)

add_text(s, Inches(0.7), Inches(5.0), Inches(12), Inches(0.4),
         "OPEN THE LAST ALERT → ALL EVENTS — SHOW:",
         size=13, bold=True, color=AMBER)
fields = [
    ("Event source", "Security"),
    ("Event ID", "1102"),
    ("Subject_username", "phantomlab"),
]
for i, (k, v) in enumerate(fields):
    y = Inches(5.45) + Inches(0.45) * i
    add_text(s, Inches(1.0), y, Inches(3.0), Inches(0.4),
             k, size=14, color=MUTED)
    add_text(s, Inches(4.2), y, Inches(7.5), Inches(0.4),
             v, size=14, color=GREEN if k == "Event ID" else WHITE, font=FONT_MONO)

add_text(s, Inches(0.7), Inches(7.0), Inches(12), Inches(0.4),
         '👉 subject_username = phantomlab — THE SAME ACCOUNT FROM THE LATERAL',
         size=13, bold=True, color=AMBER)


# ─── Slide 19: Attribution chain — one story ──────────────────────
s = prs.slides.add_slide(blank); set_slide_bg(s)
add_accent_bar(s, Inches(0), color=GREEN)
slide_title(s, "The complete attribution chain",
            "From phishing email to wiped event log — one story", eyebrow_color=GREEN)

steps = [
    ("Phishing email",       "WINWORD → cmd.exe spawn"),
    ("Implant beacon",       "mars.exe HTTP outbound (Act 1)"),
    ("Credential dump",      "rundll32 → LSASS handle (Act 2)"),
    ("Defender silenced",    "Operational 5001 captured anyway (Act 3)"),
    ("Persistence trio",     "User · Run key · Scheduled task (Act 4)"),
    ("Lateral to victim2",   "phantomlab via SMB + WinRM (Act 5)"),
    ("Stage + exfil",        "Compress-Archive → HTTP POST (Act 6)"),
    ("Anti-forensics",       "Security 1102 — wevtutil cl Security (Act 7)"),
]
top = Inches(1.8)
row_h = Inches(0.6)
for i, (k, v) in enumerate(steps):
    y = top + row_h * i
    badge_color = RED if i in (0, 2, 4, 6) else AMBER
    badge = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.9), y + Inches(0.05),
                               Inches(0.4), Inches(0.4))
    badge.fill.solid(); badge.fill.fore_color.rgb = badge_color
    badge.line.fill.background()
    tf = badge.text_frame
    tf.margin_left = tf.margin_right = Emu(0); tf.margin_top = tf.margin_bottom = Emu(0)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    r = tf.paragraphs[0].add_run(); r.text = str(i + 1)
    r.font.name = FONT_BODY; r.font.size = Pt(14); r.font.bold = True
    r.font.color.rgb = BG
    add_text(s, Inches(1.45), y + Inches(0.1), Inches(4.0), Inches(0.4),
             k, size=15, bold=True, color=WHITE)
    add_text(s, Inches(5.5), y + Inches(0.1), Inches(7.5), Inches(0.4),
             v, size=13, color=MUTED, font=FONT_MONO)

add_text(s, Inches(0.7), Inches(6.85), Inches(12), Inches(0.4),
         "\"ONE story, ONE incident, traceable end to end. THAT is what XDR is supposed to do.\"",
         size=15, bold=True, color=GREEN)


# ─── Slide 20: Response actions ────────────────────────────────────
s = prs.slides.add_slide(blank); set_slide_bg(s)
add_accent_bar(s, Inches(0), color=BLUE)
slide_title(s, "Response phase · close the loop",
            "Demonstrate · don't necessarily execute everything", eyebrow_color=BLUE)

actions = [
    ("1", "Isolate Endpoint",
     "Asset view → Initiate Response Action → Isolate",
     "Cuts the host off the network in seconds. End user still has the workstation but can't reach corp resources OR C2."),
    ("2", "Live Terminal cleanup",
     "Live Terminal → run the cleanup snippet (next slide)",
     "Now we use the tool we previewed in Act 4."),
    ("3", "Search & Destroy",
     "For staged exfil zip + mars.exe implant",
     "Wipes the artifacts across both hosts."),
    ("4", "XSOAR Playbook",
     "Run Playbook on the incident",
     "Auto-isolation + password reset + IR notification, one click."),
    ("5", "Close incident",
     "Resolution Status: Resolved Threat Handled",
     "Add notes referencing the demo's findings."),
]
top = Inches(1.7)
row_h = Inches(1.0)
for i, (n, action, where, why) in enumerate(actions):
    y = top + row_h * i
    badge = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.7), y + Inches(0.2),
                               Inches(0.5), Inches(0.5))
    badge.fill.solid(); badge.fill.fore_color.rgb = BLUE
    badge.line.fill.background()
    tf = badge.text_frame
    tf.margin_left = tf.margin_right = Emu(0); tf.margin_top = tf.margin_bottom = Emu(0)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    r = tf.paragraphs[0].add_run(); r.text = n
    r.font.name = FONT_BODY; r.font.size = Pt(20); r.font.bold = True
    r.font.color.rgb = BG
    add_text(s, Inches(1.35), y + Inches(0.1), Inches(4.5), Inches(0.4),
             action, size=15, bold=True, color=WHITE)
    add_text(s, Inches(6.0), y + Inches(0.1), Inches(7.0), Inches(0.4),
             where, size=12, color=AMBER, font=FONT_MONO)
    add_text(s, Inches(1.35), y + Inches(0.55), Inches(11.5), Inches(0.5),
             why, size=12, color=MUTED)


# ─── Slide 21: Live Terminal cleanup script ───────────────────────
s = prs.slides.add_slide(blank); set_slide_bg(s)
add_accent_bar(s, Inches(0), color=BLUE)
slide_title(s, "Live Terminal cleanup — actually run this",
            "One-shot remediation of the persistence trio", eyebrow_color=BLUE)

code_block(s, Inches(0.7), Inches(1.8), Inches(11.9), Inches(3.6), [
    "# Remove the Run key planted in Act 4",
    "Remove-ItemProperty -Path 'HKCU:\\Software\\Microsoft\\Windows\\CurrentVersion\\Run' \\",
    "    -Name PhantomUpdater",
    "",
    "# Delete the scheduled task",
    "schtasks /Delete /TN PhantomMaintenance /F",
    "",
    "# Delete the rogue local user",
    "net user T1136.001_PowerShell /Delete",
    "",
    "# Re-enable Defender real-time monitoring",
    "Set-MpPreference -DisableRealtimeMonitoring $false",
], size=13)

add_text(s, Inches(0.7), Inches(5.7), Inches(12), Inches(0.4),
         "WHILE THE OPERATOR-FACING TERMINAL RUNS:",
         size=12, bold=True, color=AMBER)
add_text(s, Inches(0.7), Inches(6.1), Inches(12), Inches(1.3),
         ("\"Notice I'm not RDP-ing in. I'm not asking the desktop team for "
          "access. This is signed PowerShell running through the XDR agent "
          "channel. Every command runs in the host's own context — but the "
          "audit trail of WHAT I ran lives in XDR's Live Terminal history "
          "for the incident report.\""),
         size=14, color=WHITE)


# ─── Slide 22: Q&A talking points ─────────────────────────────────
s = prs.slides.add_slide(blank); set_slide_bg(s)
add_accent_bar(s, Inches(0), color=AMBER)
slide_title(s, "Likely Q&A — quick answers",
            "Anticipate the doubts; have one-liners ready", eyebrow_color=AMBER)

qa = [
    ("How would you have noticed if the attacker had been quieter?",
     "Behavioral analytics on the discovery burst (Act 2) fires regardless of Defender's state. XDR's detection isn't dependent on AV remaining alive."),
    ("What if the attacker uses encrypted C2?",
     "DNS tunneling (Act 6) and volumetric outbound (Act 6) are still visible. We lose payload inspection but keep behavioral signal."),
    ("Could you have caught this earlier?",
     "The phishing email's macro execution (Act 1 step 2) is the first hard signal. With managed detection, this fires within seconds of macro spawn — pre-emptive isolate possible."),
    ("How does this scale to 10,000 endpoints?",
     "XQL hunts. The queries I showed run across the whole environment in seconds. Pin them as scheduled hunts → continuous detections."),
    ("Where's XSIAM in this picture?",
     "Same data, different prioritization layer. XDR shows you the incident; XSIAM auto-clusters incidents from this attack class + suggests playbooks."),
]
top = Inches(1.7)
row_h = Inches(1.1)
for i, (q, a) in enumerate(qa):
    y = top + row_h * i
    add_text(s, Inches(0.7), y, Inches(12), Inches(0.4),
             "Q: " + q, size=14, bold=True, color=AMBER)
    add_text(s, Inches(0.95), y + Inches(0.4), Inches(11.8), Inches(0.7),
             "→ " + a, size=12, color=WHITE)


# ─── Slide 23: XQL reference card ─────────────────────────────────
s = prs.slides.add_slide(blank); set_slide_bg(s)
add_accent_bar(s, Inches(0), color=PURPLE)
slide_title(s, "XQL reference — save these to your library before the demo",
            "Pin to Query Builder · Saved Queries · open by name", eyebrow_color=PURPLE)

queries = [
    ("#1", "Discovery burst",       "xdr_data | filter event_type=ENUM.PROCESS_START and actor in (powershell, cmd) and action in (whoami, net, systeminfo, arp, netstat, route, cmdkey)"),
    ("#2", "LSASS minidump",        "xdr_data | filter event_type=ENUM.PROCESS_ACCESS and target_image_name='lsass.exe' and actor='rundll32.exe'"),
    ("#3", "Defender disable",      "xdr_data | filter event_type=ENUM.REGISTRY and registry_path contains 'Windows Defender' and registry_value_name='DisableAntiSpyware'"),
    ("#4", "Lateral SMB+WinRM",     "authentication | filter dst_host_name='xdragent2' and src_host_name='xdragent' and action_username='phantomlab'"),
    ("#5", "DNS tunneling",         "xdr_data | filter event_type=ENUM.DNS | comp count_distinct(dns_query_name) by bin(_time,1m), parent_domain | filter unique >= 10"),
    ("#6", "Archive in TEMP",       "xdr_data | filter event_type=ENUM.FILE and action_file_extension='zip' and action_file_path contains 'Temp'"),
    ("#7", "Security 1102 backstop","xdr_data | filter event_type=ENUM.EVENT_LOG and event_id=1102"),
]
top = Inches(1.7)
row_h = Inches(0.6)
for i, (n, title, body) in enumerate(queries):
    y = top + row_h * i
    add_text(s, Inches(0.7), y, Inches(0.5), Inches(0.4),
             n, size=14, bold=True, color=PURPLE, font=FONT_MONO)
    add_text(s, Inches(1.25), y, Inches(2.8), Inches(0.4),
             title, size=13, bold=True, color=WHITE)
    add_text(s, Inches(4.15), y + Inches(0.05), Inches(8.9), Inches(0.5),
             body, size=10, color=GREEN, font=FONT_MONO)


# ─── Slide 24: Closing ─────────────────────────────────────────────
s = prs.slides.add_slide(blank); set_slide_bg(s)
add_accent_bar(s, Inches(0), color=GREEN)
add_text(s, Inches(0.7), Inches(1.0), Inches(12), Inches(0.5),
         "DEMO COMPLETE", size=14, bold=True, color=GREEN)
add_text(s, Inches(0.7), Inches(1.6), Inches(12), Inches(1.4),
         "One incident.\nOne investigator.\nOne console.", size=44, bold=True, color=WHITE)

add_panel(s, Inches(0.7), Inches(4.5), Inches(11.9), Inches(2.0), fill=PANEL)
add_text(s, Inches(1.0), Inches(4.7), Inches(11.3), Inches(1.8),
         ("Cortex XDR turned 22 attacker steps across 2 endpoints into ONE "
          "investigative story — visible end-to-end, from phishing macro to "
          "wiped event log, in about 8 minutes of analyst time.\n\n"
          "The alternative is what most SOCs still live with today: 22+ "
          "disconnected alerts across 6 tools, dispatched to different "
          "analysts, none of whom see the lateral edge."),
         size=15, color=WHITE)

add_text(s, Inches(0.7), Inches(6.85), Inches(12), Inches(0.4),
         "Sources: docs-cortex.paloaltonetworks.com · MITRE ATT&CK · Cortex XQL reference",
         size=11, color=MUTED)


# ─── Save ──────────────────────────────────────────────────────────

out = Path(__file__).resolve().parent / "case_1794_demo_deck.pptx"
prs.save(str(out))
print(f"✓ Wrote {out} ({out.stat().st_size:,} bytes, {len(prs.slides)} slides)")
